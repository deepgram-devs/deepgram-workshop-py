# /// script
# requires-python = ">=3.13"
# dependencies = ["python-pptx==1.0.2"]
# ///
"""Generate the workshop projection deck. Maintainers only.

    uv run scripts/build_slides.py

Writes slides/deepgram-voice-agent-workshop.pptx, overwriting it. The deck is
generated rather than hand-edited so a timing change in FACILITATOR.md can be
made in one place and stay true on the screen.

The deck is a visual aid for a room, not a talk. Only the Step 0 block is
presented; every other slide exists to be left on the projector while people
work, so anyone who looks up can read which step the room is on, the command to
run, and what has to be true before we move. Every slide carries speaker notes.

If you change a step's clock time, duration, or sync point in FACILITATOR.md,
change it in STEPS below too.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Palette -----------------------------------------------------------------
# Deepgram's green on near-black. Amber is reserved for one thing only: a sync
# point where the room has to regroup. Coral marks a step you must not cut.

INK = RGBColor(0x07, 0x0B, 0x10)
PANEL = RGBColor(0x11, 0x1A, 0x24)
PANEL_HI = RGBColor(0x17, 0x22, 0x2E)
CONSOLE = RGBColor(0x04, 0x07, 0x0A)
LINE = RGBColor(0x21, 0x30, 0x40)
GREEN = RGBColor(0x13, 0xEF, 0x93)
GREEN_DIM = RGBColor(0x0D, 0x6B, 0x49)
GREEN_TINT = RGBColor(0x08, 0x1F, 0x18)
AMBER = RGBColor(0xFF, 0xB4, 0x54)
AMBER_TINT = RGBColor(0x2A, 0x1F, 0x0D)
CORAL = RGBColor(0xFF, 0x7A, 0x6B)
TEXT = RGBColor(0xE9, 0xEF, 0xF6)
MUTED = RGBColor(0x93, 0xA4, 0xB7)
DIM = RGBColor(0x5D, 0x6E, 0x80)

BODY = "Calibri"
MONO = "Courier New"

W = 13.3333
H = 7.5
M = 0.7  # left and right margin
CW = W - 2 * M  # content width

# Step number, name, clock time, duration, Easy Mode lab, and the sync point.
# These mirror the run of show in FACILITATOR.md. Keep them in step.
STEPS = [
    ("1", "Setup", "0:10", "15 min", "LAB 1"),
    ("2", "Connect", "0:25", "20 min", "LAB 1"),
    ("3", "Hear the agent", "0:45", "25 min", "LAB 1"),
    ("4", "Talk to the agent", "1:10", "30 min", "LAB 1"),
    ("5", "Barge-in", "1:50", "25 min", "LAB 2"),
    ("6", "Make it yours", "2:15", "20 min", "LAB 4"),
    ("7", "Function calling", "2:35", "30 min", "LAB 5"),
    ("8", "Optimization", "3:05", "15 min", "LAB 3"),
]


# --- Primitives --------------------------------------------------------------


def new_slide(prs: Presentation) -> object:
    """Add a blank slide with the deck's background colour.

    Args:
        prs: The presentation to add to.

    Returns:
        The new slide.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = INK
    return slide


def notes(slide: object, text: str) -> None:
    """Attach speaker notes to a slide.

    Args:
        slide: The slide to annotate.
        text: The notes, as plain text.
    """
    slide.notes_slide.notes_text_frame.text = text.strip()


def box(slide: object, x: float, y: float, w: float, h: float) -> object:
    """Add a text frame with no internal padding and word wrap on.

    Args:
        slide: The slide to add to.
        x: Left edge, inches.
        y: Top edge, inches.
        w: Width, inches.
        h: Height, inches.

    Returns:
        The shape's text frame.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(
    tf: object,
    text: str,
    size: float,
    color: RGBColor,
    *,
    bold: bool = False,
    font: str = BODY,
    before: float = 0,
    after: float = 0,
    spacing: float | None = None,
    align: int = PP_ALIGN.LEFT,
    caps: bool = False,
    new: bool = True,
) -> object:
    """Append a paragraph of text to a frame.

    Args:
        tf: The text frame to append to.
        text: The text. Empty strings make deliberate vertical gaps.
        size: Font size in points.
        color: Font colour.
        bold: Whether to bold it.
        font: Typeface name.
        before: Space before the paragraph, points.
        after: Space after the paragraph, points.
        spacing: Line spacing multiple, or None for the default.
        align: Paragraph alignment.
        caps: Whether to upper-case the text.
        new: False reuses the frame's first, empty paragraph.

    Returns:
        The paragraph that was written.
    """
    p = tf.paragraphs[0] if (not new and not tf.paragraphs[0].runs) else tf.add_paragraph()
    run = p.add_run()
    run.text = text.upper() if caps else text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    return p


def block(
    slide: object,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor | None = PANEL,
    line: RGBColor | None = None,
    radius: float | None = 0.035,
) -> object:
    """Draw a filled panel, rounded unless radius is None.

    Args:
        slide: The slide to draw on.
        x: Left edge, inches.
        y: Top edge, inches.
        w: Width, inches.
        h: Height, inches.
        fill: Fill colour, or None for no fill.
        line: Outline colour, or None for no outline.
        radius: Corner radius as a fraction of the short side, or None for square corners.

    Returns:
        The shape.
    """
    shape_type = MSO_SHAPE.RECTANGLE if radius is None else MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        shape.adjustments[0] = radius
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def console(slide: object, x: float, y: float, w: float, lines: list[tuple[str, RGBColor]], size: float = 12) -> float:
    """Draw a terminal panel with one monospaced line per entry.

    Args:
        slide: The slide to draw on.
        x: Left edge, inches.
        y: Top edge, inches.
        w: Width, inches.
        lines: Text and colour for each line.
        size: Font size in points.

    Returns:
        The bottom edge of the panel, inches.
    """
    pad = 0.2
    lh = size * 1.42 / 72
    h = len(lines) * lh + pad * 2
    block(slide, x, y, w, h, fill=CONSOLE, line=LINE)
    tf = box(slide, x + pad, y + pad - 0.02, w - pad * 2, h - pad * 2)
    for i, (text, color) in enumerate(lines):
        write(tf, text, size, color, font=MONO, spacing=1.18, new=i > 0)
    return y + h


def pause_glyph(slide: object, x: float, y: float, size: float = 0.2, color: RGBColor = AMBER) -> None:
    """Draw a pause icon as two bars, rather than trusting a font to have one.

    Args:
        slide: The slide to draw on.
        x: Left edge, inches.
        y: Top edge, inches.
        size: Height of the bars, inches.
        color: Bar colour.
    """
    bar = size * 0.3
    for dx in (0, bar * 2):
        block(slide, x + dx, y, bar, size, fill=color, line=None, radius=0.4)


def heading(slide: object, title: str, kicker: str | None = None, lede: str | None = None) -> float:
    """Write a slide's kicker, title, and standfirst.

    Args:
        slide: The slide to write on.
        title: The slide title.
        kicker: A small upper-case label above the title, or None.
        lede: A line of standfirst below the title, or None.

    Returns:
        The y coordinate content can start at, inches.
    """
    y = 0.52
    if kicker:
        tf = box(slide, M, y, CW, 0.26)
        write(tf, kicker, 10.5, GREEN, bold=True, caps=True, new=False)
        y += 0.34
    tf = box(slide, M, y, CW, 0.6)
    write(tf, title, 31, TEXT, bold=True, new=False)
    y += 0.78
    if lede:
        tf = box(slide, M, y, min(CW, 10.6), 0.6)
        write(tf, lede, 14.5, MUTED, spacing=1.28, new=False)
        y += 0.42 + 0.28 * (len(lede) // 105)
    return y + 0.42


def footnote(slide: object, text: str, y: float = 6.72, color: RGBColor = DIM, size: float = 12) -> None:
    """Write a quiet line of text along the bottom of a slide.

    Args:
        slide: The slide to write on.
        text: The note.
        y: Top edge, inches.
        color: Text colour.
        size: Font size in points.
    """
    tf = box(slide, M, y, CW, 0.5)
    write(tf, text, size, color, spacing=1.25, new=False)


def rail(slide: object, current: int) -> None:
    """Draw the eight-step progress rail, so a room can read its position.

    Args:
        slide: The slide to draw on.
        current: The step number now in progress, 1 to 8. 0 draws none filled.
    """
    size, gap, y = 0.15, 0.1, 7.03
    for i in range(1, 9):
        x = M + (i - 1) * (size + gap)
        if i < current:
            block(slide, x, y, size, size, fill=GREEN_DIM, radius=0.3)
        elif i == current:
            block(slide, x, y, size, size, fill=GREEN, radius=0.3)
        else:
            block(slide, x, y, size, size, fill=None, line=LINE, radius=0.3)
    tf = box(slide, M + 8 * (size + gap) + 0.12, y - 0.04, 3.4, 0.28)
    write(tf, f"step {current} of 8", 9.5, DIM, caps=True, new=False)


def cards(
    slide: object,
    y: float,
    h: float,
    items: list[tuple[str, str]],
    *,
    label_color: RGBColor = GREEN,
    label_size: float = 13,
    body_size: float = 12.5,
    label_font: str = BODY,
    gap: float = 0.22,
    x: float = M,
    w: float = CW,
) -> None:
    """Lay out equal-width cards across the slide, each a label over a paragraph.

    Args:
        slide: The slide to draw on.
        y: Top edge, inches.
        h: Card height, inches.
        items: Label and body text for each card.
        label_color: Colour of the card labels.
        label_size: Label font size in points.
        body_size: Body font size in points.
        label_font: Typeface for the labels.
        gap: Space between cards, inches.
        x: Left edge of the row, inches.
        w: Total width of the row, inches.
    """
    n = len(items)
    cw = (w - gap * (n - 1)) / n
    for i, (label, body) in enumerate(items):
        cx = x + i * (cw + gap)
        block(slide, cx, y, cw, h, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.24, y + 0.22, cw - 0.48, h - 0.44)
        write(tf, label, label_size, label_color, bold=True, font=label_font, new=False)
        write(tf, body, body_size, MUTED, spacing=1.3, before=7)


def numbered_rows(
    slide: object,
    y: float,
    items: list[tuple[str, str, str]],
    *,
    row_h: float = 0.86,
    gap: float = 0.14,
    body_size: float = 12.5,
    head_size: float = 14,
    num_color: RGBColor = GREEN,
) -> float:
    """Draw rows of a numbered disc, a bold heading, and a description.

    Args:
        slide: The slide to draw on.
        y: Top edge of the first row, inches.
        items: Number, heading, and description for each row.
        row_h: Row height, inches.
        gap: Space between rows, inches.
        body_size: Description font size in points.
        head_size: Heading font size in points.
        num_color: Colour of the numbered discs.

    Returns:
        The bottom edge of the last row, inches.
    """
    d = 0.42
    for i, (num, head, desc) in enumerate(items):
        ry = y + i * (row_h + gap)
        disc = block(slide, M, ry + (row_h - d) / 2 - 0.16, d, d, fill=None, line=num_color, radius=0.5)
        disc.shadow.inherit = False
        tf = box(slide, M, ry + (row_h - d) / 2 - 0.09, d, 0.3)
        write(tf, num, 12.5, num_color, bold=True, align=PP_ALIGN.CENTER, new=False)
        tf = box(slide, M + d + 0.28, ry, CW - d - 0.28, row_h)
        write(tf, head, head_size, TEXT, bold=True, new=False)
        write(tf, desc, body_size, MUTED, spacing=1.28, before=4)
    return y + len(items) * (row_h + gap)


def gate_bar(slide: object, text: str, y: float = 6.06, color: RGBColor = AMBER, tint: RGBColor = AMBER_TINT) -> None:
    """Draw the sync-point bar: the one thing on a gate slide to say out loud.

    Args:
        slide: The slide to draw on.
        text: The condition that has to be true before the room moves on.
        y: Top edge, inches.
        color: Accent colour.
        tint: Panel fill.
    """
    h = 0.76
    block(slide, M, y, CW, h, fill=tint, line=color)
    pause_glyph(slide, M + 0.3, y + h / 2 - 0.11, 0.22, color)
    tf = box(slide, M + 0.95, y + 0.14, CW - 1.3, h - 0.28)
    write(tf, "before we move on", 9.5, color, bold=True, caps=True, new=False)
    write(tf, text, 14, TEXT, bold=True, before=3)


def gate_slide(
    prs: Presentation,
    num: str,
    name: str,
    clock: str,
    duration: str,
    lab: str,
    blurb: str,
    cmd: list[str],
    cmd_note: str,
    panel_label: str,
    panel_items: list[str],
    gate: str,
    *,
    tag: str | None = None,
    tag_color: RGBColor = CORAL,
    takehome: str | None = None,
) -> object:
    """Build a step gate slide, the deck's workhorse.

    Left column carries the step's identity and the command; right column carries
    what the step is actually about; the bar along the bottom carries the sync
    point. Designed to be legible from the back of a room for twenty minutes.

    Args:
        prs: The presentation to add to.
        num: Step number as it appears on screen, e.g. "01".
        name: Step name.
        clock: Clock time in the run of show, e.g. "0:10".
        duration: How long the step takes, e.g. "15 min".
        lab: The matching Easy Mode lab, e.g. "LAB 1".
        blurb: One sentence on what the step is for.
        cmd: Terminal lines to show.
        cmd_note: A quiet line under the terminal panel.
        panel_label: Heading for the right-hand panel.
        panel_items: Bulleted lines for the right-hand panel.
        gate: The sync-point condition.
        tag: An optional flag such as "DO NOT CUT".
        tag_color: Colour for that flag.
        takehome: An optional line about the matching optional step.

    Returns:
        The slide.
    """
    slide = new_slide(prs)

    tf = box(slide, M, 0.5, 6.2, 0.3)
    write(tf, f"{clock}   ·   {duration}", 13, GREEN, bold=True, caps=True, new=False)

    tf = box(slide, 7.0, 0.5, CW - 6.3, 0.3)
    write(tf, f"easy mode — {lab}", 11, DIM, bold=True, caps=True, align=PP_ALIGN.RIGHT, new=False)

    tf = box(slide, M - 0.06, 0.88, 3.0, 1.5)
    write(tf, num, 96, GREEN_DIM, bold=True, new=False)

    tf = box(slide, M, 2.28, 6.1, 0.66)
    write(tf, name, 37, TEXT, bold=True, new=False)

    if tag:
        block(slide, M, 3.06, 0.06 + len(tag) * 0.088, 0.32, fill=None, line=tag_color, radius=0.4)
        tf = box(slide, M, 3.13, 0.06 + len(tag) * 0.088, 0.24)
        write(tf, tag, 10, tag_color, bold=True, caps=True, align=PP_ALIGN.CENTER, new=False)

    by = 3.56 if tag else 3.14
    tf = box(slide, M, by, 5.9, 1.0)
    write(tf, blurb, 15, MUTED, spacing=1.3, new=False)

    cy = 4.5 if tag else 4.22
    bottom = console(slide, M, cy, 5.9, [(c, GREEN if c.startswith("uv") or c.startswith("cp") else TEXT) for c in cmd])
    if cmd_note:
        tf = box(slide, M, bottom + 0.14, 5.9, 0.4)
        write(tf, cmd_note, 11, DIM, spacing=1.25, new=False)

    px, pw = 7.0, W - 7.0 - M
    # Estimate the panel's height from its text rather than fixing it, so a
    # two-item panel doesn't leave a third of the slide empty.
    est = 0.62 + 0.26
    for item in panel_items:
        est += max(1, -(-len(item) // 48)) * 0.216 + 0.16
    ph = min(4.06 if takehome else 5.02, max(3.1, est))
    block(slide, px, 0.88, pw, ph, fill=PANEL, line=LINE)
    tf = box(slide, px + 0.3, 1.12, pw - 0.6, ph - 0.5)
    write(tf, panel_label, 11, GREEN, bold=True, caps=True, new=False)
    for i, item in enumerate(panel_items):
        write(tf, item, 13.5, TEXT if i == 0 and panel_label.startswith("THE") else MUTED, spacing=1.32, before=11)

    if takehome:
        ty = 0.88 + ph + 0.22
        block(slide, px, ty, pw, 0.7, fill=GREEN_TINT, line=GREEN_DIM)
        tf = box(slide, px + 0.3, ty + 0.15, pw - 0.6, 0.44)
        write(tf, takehome, 11.5, GREEN, bold=True, new=False)

    gate_bar(slide, gate)
    rail(slide, int(num))
    return slide


# --- The deck ----------------------------------------------------------------


def slide_title(prs: Presentation) -> None:
    """Title slide. Up while people settle in, with the command already on it.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    tf = box(slide, M, 1.16, CW, 0.3)
    write(tf, "deepgram developer workshop", 12, GREEN, bold=True, caps=True, new=False)

    tf = box(slide, M, 1.62, 8.6, 1.9)
    write(tf, "Build a Real-Time", 52, TEXT, bold=True, spacing=1.02, new=False)
    write(tf, "Voice Agent", 52, GREEN, bold=True, spacing=1.02)

    tf = box(slide, M, 3.72, 8.9, 0.9)
    write(tf, "Python Edition", 19, TEXT, bold=True, new=False)
    write(tf, "Flux speech-to-text, an LLM, and Flux TTS on one WebSocket", 15, MUTED, before=6)
    write(tf, "Roughly three hours   ·   hands-on   ·   on your own machine", 13, DIM, before=4)

    bottom = console(
        slide,
        M,
        5.28,
        6.5,
        [("uv sync", GREEN), ("cp .env.example .env", GREEN), ("uv run steps/01-setup/main.py", GREEN)],
    )
    tf = box(slide, M, bottom + 0.16, 6.5, 0.3)
    write(tf, "Start here while people settle in.", 12, DIM, new=False)

    notes(
        slide,
        """
Leave this up as the room fills. Nothing to present.

Say, once, while people are sitting down: if you have not run the setup check
yet, run it now -- that command on the screen -- because it is the only thing
that can go wrong in a way that costs you twenty minutes later.

The single biggest predictor of whether this session goes well is how many
people arrive with a working environment. If you sent the pre-event email a
week out, most of the room is fine. Find the ones who aren't now, not at 0:10.

Wired headphones. Ask out loud who doesn't have any and pair them with someone
who does, or seat them where the room noise is lowest. Step 5 is where the
difference shows.
        """,
    )


def slide_what_youll_build(prs: Presentation) -> None:
    """The promise, with the four numbers that make it concrete.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "What you'll build",
        lede=(
            "A voice agent you can hold a conversation with. It listens on Flux, thinks with an LLM, answers in a "
            "natural voice, stops talking the moment you interrupt it, and calls your Python when it needs "
            "something it can't know."
        ),
    )

    stats = [
        ("~350", "lines of Python", "when you're done"),
        ("8", "steps", "every one a complete running program"),
        ("1", "API key", "covers all three models"),
        ("<$1", "of credit", "spent per attendee"),
    ]
    gap, n = 0.22, len(stats)
    cw = (CW - gap * (n - 1)) / n
    for i, (value, label, sub) in enumerate(stats):
        cx = M + i * (cw + gap)
        block(slide, cx, y + 0.1, cw, 2.2, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.26, y + 0.4, cw - 0.52, 0.9)
        write(tf, value, 46, GREEN, bold=True, new=False)
        tf = box(slide, cx + 0.26, y + 1.34, cw - 0.52, 0.8)
        write(tf, label, 14.5, TEXT, bold=True, new=False)
        write(tf, sub, 12, MUTED, spacing=1.25, before=4)

    block(slide, M, y + 2.62, CW, 0.92, fill=GREEN_TINT, line=GREEN_DIM)
    tf = box(slide, M + 0.34, y + 2.84, CW - 0.68, 0.5)
    write(tf, "The agent is entirely Python.", 15, GREEN, bold=True, new=False)
    write(
        tf,
        "The browser is only the microphone and the speaker — and swapping it for a Twilio media stream never "
        "touches a line of agent code.",
        13.5,
        MUTED,
        before=4,
    )

    notes(
        slide,
        """
Thirty seconds on this slide. It is a promise, not a lesson.

The number that matters to the room is 1 API key. Your Deepgram key pays for
speech-to-text, the LLM, and text-to-speech alike -- Deepgram brokers the LLM
call, so there is no OpenAI account, no second key, and no second bill. People
who have wired up a voice pipeline before will not expect that.

New accounts get $200 in credit and a full three-hour run costs well under a
dollar per person, so nobody needs to worry about the meter.

If someone asks about the browser: it is there for echo cancellation and for
the microphone permission model, not because the agent needs a front end. That
choice is what makes WSL work and what stops Step 5 turning into twenty minutes
of debugging a laptop speaker.
        """,
    )


def slide_prerequisites(prs: Presentation) -> None:
    """What has to be true of each machine before Step 1.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "What you need on your machine",
        lede="One dependency to install, one key to paste, and headphones. That's the whole bar.",
    )

    numbered_rows(
        slide,
        y - 0.08,
        [
            (
                "1",
                "uv",
                "The only thing to install. It fetches Python 3.13 and every dependency itself, so you don't need "
                "Python already. One virtual environment serves every step — you run uv sync once.",
            ),
            (
                "2",
                "A free Deepgram API key",
                "console.deepgram.com/signup — signup takes a minute, and new accounts get $200 in credit.",
            ),
            (
                "3",
                "A current browser",
                "Chrome, Firefox, or Safari. The browser is the microphone and the speaker, so it needs "
                "getUserMedia and AudioWorklet.",
            ),
            (
                "4",
                "Wired headphones",
                "Your browser cancels most of the echo from your speakers, but not all of it on every browser. "
                "Step 5 is where the difference shows.",
            ),
        ],
        row_h=0.82,
        gap=0.12,
    )

    block(slide, M, 6.1, CW, 0.82, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + 0.34, 6.28, CW - 0.68, 0.5)
    write(tf, "Locked-down laptop, Chromebook, or tablet?", 13.5, GREEN, bold=True, new=False)
    write(
        tf,
        "Easy Mode runs the same workshop in the browser, no install. Or .devcontainer/ runs the whole thing in "
        "GitHub Codespaces, with nothing local at all.",
        12.5,
        MUTED,
        before=4,
    )

    notes(
        slide,
        """
Leave this up during setup. Do not read it out.

The two you will actually be asked about:

uv -- people assume they need Python first. They don't. uv fetches Python 3.13
itself. If someone has an old Python and is worried, that is fine, uv ignores it.

Headphones -- the request in the pre-event email was not decoration. Without
them the agent can hear itself through the laptop speakers and interrupt itself,
which looks exactly like a broken barge-in implementation. That costs people
twenty minutes in Step 5 and it costs you the best moment in the workshop.

A terminal on macOS, Linux, Windows, or WSL all work. WSL is no longer a problem
because the audio is Windows' problem now, not yours -- the browser handles it.
        """,
    )


def slide_two_tracks(prs: Presentation) -> None:
    """Code track and Easy Mode, and how a mixed room stays in sync.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Two tracks, one room",
        lede="Both tracks talk to the same Voice Agent API. Nobody is doing a lesser version of this.",
    )

    tracks = [
        (
            "CODE TRACK",
            "You write the agent in Python",
            "Eight steps, each a complete runnable program. Work the TODO blocks in main.py; the next folder is "
            "the answer key.",
            "uv   ·   a terminal   ·   a browser",
            GREEN,
        ),
        (
            "EASY MODE",
            "You configure the agent in a browser",
            "Six labs in the Deepgram Playground. No install, no terminal, no key pasted into a file. You end with "
            "a config you can hand to a developer.",
            "a browser   ·   headphones",
            AMBER,
        ),
    ]
    gap = 0.3
    cw = (CW - gap) / 2
    for i, (kicker, title, body, needs, color) in enumerate(tracks):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 2.66, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.32, y + 0.28, cw - 0.64, 2.1)
        write(tf, kicker, 11, color, bold=True, caps=True, new=False)
        write(tf, title, 18, TEXT, bold=True, before=8, spacing=1.14)
        write(tf, body, 13, MUTED, before=8, spacing=1.3)
        write(tf, needs, 11.5, color, bold=True, before=12, font=MONO)

    my = y + 3.0
    tf = box(slide, M, my, CW, 0.3)
    write(tf, "how the labs line up", 10.5, DIM, bold=True, caps=True, new=False)
    pairs = [
        ("Steps 1–4", "Lab 1"),
        ("Step 5", "Lab 2"),
        ("Step 6", "Lab 4"),
        ("Step 7", "Lab 5"),
        ("Step 8", "Lab 3"),
        ("Finished", "Lab 6"),
    ]
    g, n = 0.16, len(pairs)
    pw = (CW - g * (n - 1)) / n
    for i, (step, lab) in enumerate(pairs):
        px = M + i * (pw + g)
        block(slide, px, my + 0.42, pw, 0.86, fill=PANEL_HI, line=LINE)
        tf = box(slide, px + 0.16, my + 0.58, pw - 0.32, 0.6)
        write(tf, step, 13, TEXT, bold=True, align=PP_ALIGN.CENTER, new=False)
        write(tf, lab, 12, AMBER, bold=True, align=PP_ALIGN.CENTER, before=3)

    footnote(
        slide,
        "Easy Mode reaches turn detection earlier than the code track does, which is the only reason its lab "
        "numbers don't run in order. Every checkpoint still lands in the same place.",
    )

    notes(
        slide,
        """
Worth thirty seconds, because it decides whether anyone feels stranded.

Say out loud: if your environment went sideways, or you did not come here to
write Python, Easy Mode is not a consolation prize. It covers every concept in
this workshop through the Playground, and you will finish each lab faster than
the code track finishes its step. Spend the spare minutes on the Going further
prompts in the lab.

Also say: the Playground's settings panel is a form over a JSON document, and
the developers next to you are typing that same document by hand. Once someone
sees that, the two tracks stop looking like different activities.

Pair people up if the room allows it. An Easy Mode attendee sitting next to a
code track attendee is the best version of this -- at Step 5 the Easy Mode
person gets barge-in for free and can watch their neighbour implement it.
        """,
    )


def slide_run_of_show(prs: Presentation) -> None:
    """The three-hour map, with the sync points marked.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Run of show",
        lede=(
            "Three hours, one break. Amber marks where we regroup — nobody moves past a gate until the room is "
            "together."
        ),
    )

    rows = [
        ("0:00", "Overview", "The parts of a voice agent", False),
        ("0:10", "Setup", "Everyone runs the checker", True),
        ("0:25", "Connect", "The settings handshake", True),
        ("0:45", "Hear the agent", "Playback — the greeting plays", True),
        ("1:10", "Talk to the agent", "Microphone in, a conversation", True),
        ("1:40", "Break", "Ten minutes", False),
        ("1:50", "Barge-in", "Interrupting it mid-sentence", True),
        ("2:15", "Make it yours", "Prompt, persona, voice", True),
        ("2:35", "Function calling", "It answers from your data", True),
        ("3:05", "Optimization", "Turn thresholds and latency", False),
        ("3:20", "Wrap", "Where to go next, and questions", False),
    ]

    gap = 0.34
    cw = (CW - gap) / 2
    per = 6
    for i, (clock, name, what, gated) in enumerate(rows):
        col, idx = (0, i) if i < per else (1, i - per)
        rx = M + col * (cw + gap)
        ry = y + idx * 0.68
        block(slide, rx, ry, cw, 0.58, fill=PANEL if gated else None, line=LINE)
        tf = box(slide, rx + 0.22, ry + 0.16, 0.72, 0.3)
        write(tf, clock, 13, GREEN if gated else DIM, bold=True, font=MONO, new=False)
        tf = box(slide, rx + 1.06, ry + 0.14, cw - 1.7, 0.34)
        write(tf, name, 14, TEXT, bold=True, new=False)
        tf = box(slide, rx + 1.06, ry + 0.14, cw - 1.06 - (0.62 if gated else 0.24), 0.34)
        write(tf, what, 11.5, MUTED, align=PP_ALIGN.RIGHT, new=False)
        if gated:
            pause_glyph(slide, rx + cw - 0.34, ry + 0.19, 0.19)

    block(slide, M + cw + gap, y + 5 * 0.68, cw, 1.32, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + cw + gap + 0.26, y + 5 * 0.68 + 0.2, cw - 0.52, 1.0)
    write(tf, "if we run long", 10.5, CORAL, bold=True, caps=True, new=False)
    write(
        tf,
        "Step 8 goes first — it's dials rather than code and makes a clean take-home. Step 7's exercises go second. "
        "Step 5 never goes.",
        12.5,
        MUTED,
        before=6,
        spacing=1.28,
    )

    notes(
        slide,
        """
Show it, name the break, move on. Do not walk the room through eleven rows.

The two sentences worth saying:

"The amber rows are where we stop and make sure everyone is together. If you're
behind at one of those, say so -- that's what they're for."

"Every step folder contains a complete, working version of the step before it.
So if you fall behind, you skip ahead and keep going. You lose the typing, not
the workshop." Say this now, at the first gate, and again at the break. If you
don't say it, people fall behind quietly and disengage.

Timings come from FACILITATOR.md and assume a 3-hour slot. A 90-minute run is
Steps 0 to 5 and still ends with a working voice agent you can interrupt.

Steps 6b and 7b are deliberately not on this map. Mention each in one sentence
when you finish the step before it, and point at the appendix.
        """,
    )


def slide_three_models(prs: Presentation) -> None:
    """Step 0: the pipeline, as three boxes and two arrows.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "A voice agent is three models in a loop",
        kicker="step 0  ·  the only block where you talk",
        lede="Audio comes in, audio goes out, and a model in the middle decides what to say.",
    )

    stages = [
        ("SPEECH TO TEXT", "Flux", "hears you", GREEN),
        ("THE BRAIN", "gpt-4o-mini", "decides what to say", TEXT),
        ("TEXT TO SPEECH", "Flux TTS", "says it", GREEN),
    ]
    gap, n = 0.62, len(stages)
    cw = (CW - gap * (n - 1)) / n
    for i, (kicker, model, role, color) in enumerate(stages):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 1.86, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.28, y + 0.26, cw - 0.56, 1.4)
        write(tf, kicker, 10.5, DIM, bold=True, caps=True, new=False)
        write(tf, model, 25, color, bold=True, before=8, font=MONO)
        write(tf, role, 14, MUTED, before=8)
        if i < n - 1:
            ax = cx + cw + 0.11
            block(slide, ax, y + 0.9, gap - 0.22, 0.045, fill=GREEN_DIM, radius=None)
            head = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ax + gap - 0.3), Inches(y + 0.79), Inches(0.17), Inches(0.26)
            )
            head.rotation = 90
            head.fill.solid()
            head.fill.fore_color.rgb = GREEN_DIM
            head.line.fill.background()
            head.shadow.inherit = False

    tf = box(slide, M, y + 2.04, cw, 0.3)
    write(tf, "your microphone", 11, DIM, caps=True, new=False)
    tf = box(slide, M + 2 * (cw + gap), y + 2.04, cw, 0.3)
    write(tf, "your speaker", 11, DIM, caps=True, align=PP_ALIGN.RIGHT, new=False)

    block(slide, M, y + 2.52, CW, 1.2, fill=GREEN_TINT, line=GREEN_DIM)
    tf = box(slide, M + 0.34, y + 2.72, CW - 0.68, 0.86)
    write(tf, "One WebSocket carries all of it.", 16, GREEN, bold=True, new=False)
    write(
        tf,
        "You could wire the three together yourself — your STT call, piped into your LLM call, piped into your TTS "
        "call. You'd spend your time on plumbing and latency, and very little on what the agent actually does.",
        13,
        MUTED,
        before=5,
        spacing=1.28,
    )

    notes(
        slide,
        """
This is the first of six slides you actually present. Ten minutes for the whole
block -- they came to code, not to watch slides.

Three models, three jobs:

Flux, speech-to-text. Deepgram's conversational model. It streams transcripts as
you speak, so nothing waits for a complete utterance, and it scores end-of-turn
confidence while the audio is still arriving. That second half is the slide after
next and it is the important one.

gpt-4o-mini, the brain. Takes the transcript, mixes it with the conversation
history and the agent's standing instructions, decides what to say back. Chosen
because it is fast and cheap, and speed matters more than raw capability when
somebody is waiting to hear a reply. The same pattern works with Anthropic,
Google, Groq, and Bedrock -- you swap a provider, not your architecture.

Flux TTS, text-to-speech. Streaming voice engine built for conversation.

Check yourself, ask the room: what are the three building blocks, and which one
holds the conversation history? Answer: STT, LLM, TTS -- and the LLM holds it.

If someone asks whether the speech models must run on Deepgram's infrastructure:
no. Flux, Nova-3, and Aura-2 are on AWS Marketplace and deploy to SageMaker
endpoints in your own VPC. That is a different architecture, not a setting, and
the Voice Agent API itself does not run there -- SageMaker's network isolation
blocks the outbound LLM calls the orchestrator makes. Park it; it's on a slide
in the appendix path if they want more.
        """,
    )


def slide_architecture(prs: Presentation) -> None:
    """Step 0: where each piece runs, and the seven hops between them.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Where the pieces actually run",
        kicker="step 0",
        lede="The agent runs on your machine. Deepgram runs the real-time speech path. Nothing is deployed anywhere.",
    )

    zones = [
        (
            "YOUR MACHINE",
            GREEN,
            [
                ("Web browser", "Microphone and speaker, echo cancellation\ngetUserMedia + AudioWorklet"),
                ("Local Python agent", "FastAPI audio bridge, your function code\n127.0.0.1:8000"),
            ],
        ),
        (
            "DEEPGRAM VOICE AGENT API",
            GREEN,
            [
                ("Agent orchestrator", "One socket carries listen, think,\nspeak, and function calls"),
                ("Flux STT  +  Flux TTS", "Turn detection server-side,\nsynthesis as the reply arrives"),
            ],
        ),
        (
            "YOUR AWS ACCOUNT  ·  STEP 6b ONLY",
            AMBER,
            [
                ("Amazon Bedrock Runtime", "Your model, your bill.\nCalled as you, with your credentials"),
            ],
        ),
    ]

    gap = 0.28
    widths = [4.24, 4.24, CW - 2 * 4.24 - 2 * gap]
    x = M
    for (label, color, items), zw in zip(zones, widths, strict=True):
        block(slide, x, y, zw, 3.1, fill=None, line=color if color is AMBER else LINE)
        tf = box(slide, x + 0.22, y + 0.18, zw - 0.44, 0.26)
        write(tf, label, 9.5, color, bold=True, caps=True, new=False)
        card_h = 1.1 if len(items) > 1 else 1.62
        for j, (title, sub) in enumerate(items):
            iy = y + 0.56 + j * 1.26
            block(slide, x + 0.2, iy, zw - 0.4, card_h, fill=PANEL, line=LINE)
            tf = box(slide, x + 0.42, iy + 0.18, zw - 0.84, card_h - 0.3)
            write(tf, title, 13.5, TEXT, bold=True, new=False)
            write(tf, sub, 11, MUTED, before=4, spacing=1.22)
        x += zw + gap

    hops = [
        ("1", "mic in"),
        ("2", "audio up"),
        ("3", "transcript + turn end"),
        ("4", "completion, signed as you"),
        ("5", "tokens back"),
        ("6", "synthesis"),
        ("7", "audio down"),
    ]
    hy = y + 3.36
    g, n = 0.12, len(hops)
    hw = (CW - g * (n - 1)) / n
    for i, (num, label) in enumerate(hops):
        hx = M + i * (hw + g)
        color = AMBER if num in {"4", "5"} else GREEN
        block(slide, hx, hy, hw, 0.68, fill=PANEL_HI, line=LINE)
        tf = box(slide, hx + 0.14, hy + 0.1, hw - 0.28, 0.24)
        write(tf, num, 11, color, bold=True, new=False)
        tf = box(slide, hx + 0.14, hy + 0.33, hw - 0.28, 0.3)
        write(tf, label, 10.5, MUTED, spacing=1.14, new=False)

    footnote(
        slide,
        "Hops 4 and 5 are Deepgram's brokered LLM in Steps 1 through 8. Only in optional Step 6b do they leave for "
        "your own AWS account.",
        y=hy + 0.86,
    )

    notes(
        slide,
        """
New slide, and the one that answers "wait, what is running where?" before
anybody has to ask it. Sixty seconds.

The shape to convey: your machine holds the browser and a small Python process.
Deepgram holds the pipeline. There is no server to deploy, nothing always-on,
and nothing accruing cost between sessions -- everything is either a local
process you started or a metered API call.

The seven hops, if you want to walk them: browser captures audio and streams it
to the local bridge over 127.0.0.1; Python forwards it to Deepgram on one
authenticated WebSocket; Flux transcribes and decides the turn ended; the
orchestrator calls the think provider; the model streams tokens back; TTS
synthesises the reply as it arrives rather than waiting for a full sentence;
agent audio returns on the same socket and plays in the browser.

Two details people ask about:

Why 127.0.0.1 and not a LAN address -- browsers only grant microphone access in
a secure context, and localhost qualifies while a LAN address does not. This is
environment failure number two later today, so planting it now pays off.

Where the credential travels in Step 6b -- your AWS access key goes into the
Settings message, over the WebSocket, to Deepgram. That is what "called as you"
means, and it's why the step asks for an IAM user scoped to one model ARN rather
than admin keys. Appendix slide has the detail.
        """,
    )


def slide_orchestration(prs: Presentation) -> None:
    """Step 0: what you'd own yourself versus what the API owns.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(slide, "The API does the orchestration", kicker="step 0")

    columns = [
        (
            "WIRE IT YOURSELF",
            DIM,
            None,
            [
                "Three services, and all the glue between them",
                "You own the handoffs and the buffering",
                "You own every millisecond between hops",
                "You own turn-taking, in your client",
            ],
        ),
        (
            "VOICE AGENT API",
            GREEN,
            GREEN_DIM,
            [
                "One WebSocket, one settings message",
                "Deepgram runs the pipeline server-side",
                "Audio bytes in, audio bytes out",
                "Turn-taking happens inside the model",
            ],
        ),
    ]
    gap = 0.36
    cw = (CW - gap) / 2
    for i, (label, color, outline, items) in enumerate(columns):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 2.9, fill=GREEN_TINT if outline else PANEL, line=outline or LINE)
        tf = box(slide, cx + 0.34, y + 0.28, cw - 0.68, 2.3)
        write(tf, label, 12, color, bold=True, caps=True, new=False)
        for item in items:
            write(tf, item, 14, TEXT if outline else MUTED, before=13, spacing=1.26)

    block(slide, M, y + 3.24, CW, 1.12, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + 0.34, y + 3.46, CW - 0.68, 0.7)
    write(tf, "You describe what you want, and you get a socket.", 16, TEXT, bold=True, new=False)
    write(
        tf,
        "Which STT model, which LLM provider and model, which TTS voice, and the agent's personality and greeting. "
        "Over the next eight steps you configure that orchestrator and react to what it sends you. You will never "
        "reimplement it.",
        13,
        MUTED,
        before=5,
        spacing=1.28,
    )

    notes(
        slide,
        """
Forty-five seconds. The left column is what most of the room expects to be
signing up for; the right column is the actual deal.

The framing that lands: you are not building a pipeline today, you are
configuring one and reacting to what it tells you. Everything you write in the
next three hours is either a settings object or an event handler.

Check yourself, ask the room: name two things the Voice Agent API handles that
you would otherwise write yourself. Answers: turn-taking, provider handoffs,
buffering, interruption signalling. Any two.

If someone pushes on lock-in or wanting their own model: good instinct, and
Step 6b is the answer -- think.endpoint plus provider credentials points the
middle of the pipeline anywhere that speaks OpenAI Chat Completions. One
sentence, then move; the appendix carries it.
        """,
    )


def slide_what_flux_changes(prs: Presentation) -> None:
    """Step 0: turn detection moves inside the model.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "What Flux changes",
        kicker="step 0  ·  slow down here",
        lede="This is where voice agents differ most from what you'd expect.",
    )

    columns = [
        (
            "TRADITIONAL PIPELINE",
            DIM,
            None,
            "Your client decides when the user has stopped talking.",
            "You measure audio energy. You set a silence threshold. You write voice activity detection — and you "
            "get it slightly wrong forever. Turn-taking lives in your code.",
        ),
        (
            "WITH FLUX",
            GREEN,
            GREEN_DIM,
            "Turn detection happens inside the model.",
            "Flux scores every turn for end-of-turn confidence as the audio streams in, and tells you what it "
            "concluded. You tune it with a number instead of maintaining a heuristic.",
        ),
    ]
    gap = 0.36
    cw = (CW - gap) / 2
    for i, (label, color, outline, head, body) in enumerate(columns):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 2.5, fill=GREEN_TINT if outline else PANEL, line=outline or LINE)
        tf = box(slide, cx + 0.34, y + 0.28, cw - 0.68, 1.9)
        write(tf, label, 12, color, bold=True, caps=True, new=False)
        write(tf, head, 16.5, TEXT, bold=True, before=9, spacing=1.16)
        write(tf, body, 13, MUTED, before=8, spacing=1.3)

    block(slide, M, y + 2.86, CW, 1.16, fill=CONSOLE, line=GREEN_DIM)
    tf = box(slide, M + 0.4, y + 3.16, CW - 0.8, 0.6)
    write(
        tf,
        "You will not write a silence threshold anywhere in this workshop.",
        22,
        GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        new=False,
    )

    notes(
        slide,
        """
Slow down. Ninety seconds, and it is the most valuable ninety seconds in the
block, because it is the thing experienced people have to unlearn.

Anyone in the room who has built a voice pipeline has written voice activity
detection, and has an opinion about it. Name that out loud -- it earns you the
next slide. Ask who has shipped a silence threshold and watch the hands.

The point is not that Flux is better at guessing. The point is that turn
detection is a modelling problem being solved by the model that already has the
audio, rather than an energy heuristic in your client that has strictly less
information to work with.

What it costs: a number. eot_threshold, 0.5 to 0.9, and that is Step 8.

What it does not remove is on the next slide, and it is the one obligation you
still own. Do not let anyone leave this block thinking barge-in is free.
        """,
    )


def slide_one_job(prs: Presentation) -> None:
    """Step 0: the single client-side obligation Flux leaves behind.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)

    tf = box(slide, M, 0.86, CW, 0.3)
    write(tf, "step 0  ·  what that leaves you", 11, GREEN, bold=True, caps=True, new=False)

    tf = box(slide, M, 1.34, 11.5, 2.1)
    write(
        tf,
        "When Flux says the user started talking,",
        38,
        TEXT,
        bold=True,
        spacing=1.1,
        new=False,
    )
    write(tf, "stop your speaker immediately.", 38, GREEN, bold=True, spacing=1.1)

    tf = box(slide, M, 3.72, 10.2, 1.4)
    write(
        tf,
        "The agent stops generating the moment it knows. But whatever it already sent is sitting in your playback "
        "buffer — easily a second or two of speech — and it will keep talking over the user until you clear it.",
        16,
        MUTED,
        spacing=1.34,
        new=False,
    )

    block(slide, M, 5.42, CW, 1.02, fill=AMBER_TINT, line=AMBER)
    pause_glyph(slide, M + 0.34, 5.79, 0.26)
    tf = box(slide, M + 1.02, 5.62, CW - 1.4, 0.64)
    write(tf, "That single responsibility gets its own step.", 15.5, TEXT, bold=True, new=False)
    write(
        tf,
        "Step 5, at 1:50. It is the difference between a demo and something a person would willingly talk to, and "
        "it is the one step we never cut.",
        13,
        AMBER,
        before=4,
        spacing=1.26,
    )

    notes(
        slide,
        """
One statement, said once, then stop talking. Let it sit for a beat.

This is the setup for the best moment in the workshop. At 1:50 they will run an
agent that ignores them mid-sentence, and because you said this at 0:07 they
will already know why.

Check yourself, ask the room: where does turn detection happen with Flux, and
what is the client still responsible for? Answer: inside the model, server-side
-- and the client is still responsible for clearing queued playback on barge-in.

Do not explain the two-queue problem now. It is Step 5's slide and it lands far
better after they have heard the bug with their own ears.

You are about six or seven minutes into the ten. One more slide, then they code.
        """,
    )


def slide_layout(prs: Presentation) -> None:
    """Step 0: how the repo is arranged, and why falling behind is survivable.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "How the workshop is laid out",
        kicker="step 0",
        lede="Every step folder holds the finished state of the step before it. Two useful things follow from that.",
    )

    cards(
        slide,
        y,
        2.02,
        [
            (
                "Every step is a complete program",
                "steps/03-hear-the-agent/main.py already does everything Steps 1 and 2 built. You open it, run it, "
                "and add the piece that step is about.",
            ),
            (
                "The next folder is the answer key",
                "Stuck on a detail in Step 3? steps/04-talk-to-the-agent/main.py is Step 3, finished. Read it "
                "rather than waiting for a helper.",
            ),
            (
                "Falling behind is survivable",
                "On Step 3 when the room hits Step 5? Skip to steps/05-barge-in/ and keep going. You lose the "
                "typing, not the workshop.",
            ),
        ],
        label_size=15,
        body_size=13,
    )

    bottom = console(
        slide,
        M,
        y + 2.34,
        CW,
        [
            ("    # TODO (Step 3.1): route the agent's audio to your speaker", DIM),
            ("    #: player.send() hands raw PCM to the playback queue.", GREEN),
            ("    # player.send(audio)", MUTED),
        ],
        size=12.5,
    )

    footnote(
        slide,
        "Work the TODO (Step N.x) blocks in order. Inside a block, #: marks the instructions — everything else is "
        "code, commented out at the indentation it belongs at. Select those lines and press Cmd+/ (Ctrl+/ on "
        "Windows and Linux) to uncomment them where they sit.",
        y=bottom + 0.2,
        size=12.5,
        color=MUTED,
    )

    notes(
        slide,
        """
Last slide of the presented block. Sixty seconds, then get them coding.

The mechanic to demonstrate rather than describe: open a main.py, show one TODO
block, select the commented lines, hit Cmd+/ and watch them uncomment at the
right indentation. Ten seconds of screen share saves fifteen questions.

Say the falling-behind rule out loud now, and say it again at the first gate:
if you are behind, skip to the next folder. Everything before it is already
done for you. People will not believe this is allowed unless you say it.

Check yourself, ask the room: you're stuck halfway through Step 5, where do you
look for the working version? Answer: the next folder.

Then: "That's the last slide I talk at you. Run Step 1."
        """,
    )


# --- Steps 1 to 5: the core --------------------------------------------------


def slide_step1(prs: Presentation) -> None:
    """Step 1 gate: prove the machine works before anything depends on it.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "01",
        "Setup",
        "0:10",
        "15 min",
        "LAB 1",
        "Prove this machine can reach Deepgram and move audio in both directions, before any of it matters.",
        ["uv sync", "cp .env.example .env", "uv run steps/01-setup/main.py"],
        "One virtual environment serves every step. You run uv sync once.",
        "WHAT IT ACTUALLY CHECKS",
        [
            "The key — with a real authenticated call, not a check that the string is non-empty",
            "The region — and whether that endpoint serves all three models",
            "Microphone permission, and your actual input level",
            "The audio APIs, and whether the browser granted echo cancellation",
            "Speaker output, with a test tone",
        ],
        "Nobody proceeds without a green key line and a green Agent started.",
    )
    notes(
        slide,
        """
Hold the room here. This gate is the one that pays for itself.

Leave the slide up. People who arrived with a working environment will be done in
two minutes; use the remaining thirteen on everyone else. One floating helper per
fifteen attendees, and the helper should be walking, not sitting.

The check finishes in a browser page, so watch for people who read the green
terminal lines and stop. Both halves have to be green.

Failures you will hit here, in order of frequency: microphone permission clicked
past; the page open on a LAN address instead of 127.0.0.1; a Bluetooth headset
that flips to a mono profile; a truncated API key. All four are on the appendix
slide -- jump there rather than describing them from memory.

Say the falling-behind rule out loud at this gate. "If you get stuck, skip to the
next folder -- it already contains a finished version of everything before it.
You lose the typing, not the workshop." This is the moment people decide whether
it is safe to admit they are behind.

Check yourself, ask the room: how many API keys does this workshop need?
Answer: one. Deepgram brokers the OpenAI call, so your Deepgram key covers the LLM.
        """,
    )


def slide_step1_working(prs: Presentation) -> None:
    """Step 1: the exact output that counts as working, on both sides.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        'Step 1 — what "working" looks like',
        lede="Both halves have to be green. The terminal checks the account; the browser page checks the machine.",
    )

    gap = 0.34
    cw = (CW - gap) / 2

    tf = box(slide, M, y - 0.06, cw, 0.3)
    write(tf, "in the terminal", 10.5, GREEN, bold=True, caps=True, new=False)
    console(
        slide,
        M,
        y + 0.3,
        cw,
        [
            ("[  OK  ] DEEPGRAM_API_KEY found (84de...c2db)", GREEN),
            ("[  OK  ] Deepgram accepted the key", GREEN),
            ("[  OK  ] Region: global", GREEN),
            ("[  OK  ] Agent started", GREEN),
            ("         flux-general-en + gpt-4o-mini", MUTED),
            ("         + flux-alexis-en", MUTED),
            ("", MUTED),
            ("Open http://127.0.0.1:8000", TEXT),
        ],
        size=12.5,
    )

    tf = box(slide, M + cw + gap, y - 0.06, cw, 0.3)
    write(tf, "in the browser page", 10.5, GREEN, bold=True, caps=True, new=False)
    console(
        slide,
        M + cw + gap,
        y + 0.3,
        cw,
        [
            ("OK   Secure context", GREEN),
            ("OK   Audio support", GREEN),
            ("OK   Deepgram API key", GREEN),
            ("OK   Microphone      (peak 0.34)", GREEN),
            ("OK   Echo cancellation granted", GREEN),
            ("OK   Speaker         tone played", GREEN),
            ("", MUTED),
            ("Ready for Step 2.", TEXT),
        ],
        size=12.5,
    )

    block(slide, M, 5.6, CW, 1.24, fill=AMBER_TINT, line=AMBER)
    tf = box(slide, M + 0.34, 5.8, CW - 0.68, 0.9)
    write(tf, '"The agent refused these settings" is not a broken machine.', 15, AMBER, bold=True, new=False)
    write(
        tf,
        "A model isn't served where that person is connecting, and Step 1 names which one. Same fix for the whole "
        "room: DEEPGRAM_REGION=global, or change the named model everywhere. Appendix has the detail.",
        13,
        MUTED,
        before=5,
        spacing=1.28,
    )

    notes(
        slide,
        """
This is a reference slide, not a talk. Leave it up beside the gate slide if you
have a second screen.

Its real job is to stop people declaring victory early. The terminal going green
proves the key and the account. It proves nothing about the microphone. The
browser page is the half that fails on the day.

Peak level on the microphone line matters: a peak of 0.00 with permission
granted usually means the wrong input device, not a broken browser.

Echo cancellation granted is worth reading when someone later reports the agent
interrupting itself. If it says the browser did not grant it, headphones are not
optional for that person.

Regional note if you are running outside global: Step 1 opens the same socket
Step 2 will, with the same three models, and reports whether the server accepted
them. That -- not the key, and not whether the host answers -- is the question
worth asking before a room of people hits Step 2.
        """,
    )


def slide_step2(prs: Presentation) -> None:
    """Step 2 gate: the handshake, and the one object that configures everything.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "02",
        "Connect",
        "0:25",
        "20 min",
        "LAB 1",
        "Open a WebSocket, describe the agent you want, and confirm the server accepted it.",
        ["uv run steps/02-connect/main.py"],
        "No audio yet. This step ends with a handshake, not a conversation.",
        "ONE OBJECT CONFIGURES ALL THREE",
        [
            "listen — Flux. Turn detection lives inside the model.",
            "think — the LLM. Brokered by Deepgram, so your key covers it.",
            "speak — Flux TTS. The flux- prefix routes to v2 Speak.",
            "greeting — what it says first. Leave it out and the agent waits for the user, which is right more "
            "often than people expect.",
        ],
        "Everyone sees >> Settings applied.",
    )
    notes(
        slide,
        """
Twenty minutes, and it is mostly typing a settings object. Let them work.

The one thing worth saying out loud: listen configures speech-to-text, and the
LLM is named under think. People reach for listen when they want to change the
model and then wonder why nothing happened.

Watch for anyone who has not seen ">> Settings applied" by about 0:40. Nine
times in ten it is a typo in a provider type string, and the server's error frame
names the field.

Check yourself, two of them:
- Which part of SETTINGS configures speech-to-text, and where is the LLM named?
  Answer: listen configures STT; the LLM is under think.
- What happens to audio you send before the handshake completes?
  Answer: the agent discards any media that arrives before SettingsApplied, so
  audio sent early is silently lost. That is why the bridge waits for the
  browser's start message before it opens the Deepgram socket at all.

Next slide is the ordering. It is worth showing even though nobody has to write
it, because every later step leans on that sequence.
        """,
    )


def slide_step2_ordering(prs: Presentation) -> None:
    """Step 2: the six-line sequence every later step depends on.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "The ordering that everything else depends on",
        lede="Inside web/session.py, _run() is six lines long. You never edit it, and every later step leans on it.",
    )

    seq = [
        ("1", "open the socket", ""),
        ("2", "register handlers", "before listening, or the first events go nowhere"),
        ("3", "start the listener", "the receive loop blocks, so it gets its own thread"),
        ("4", "send SETTINGS", ""),
        ("5", "wait for SettingsApplied", "up to ten seconds — it is not instant"),
        ("6", "stream media", "and not one byte before"),
    ]
    gap, n = 0.16, len(seq)
    cw = (CW - gap * (n - 1)) / n
    for i, (num, label, sub) in enumerate(seq):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 2.44, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.18, y + 0.2, cw - 0.36, 0.5)
        write(tf, num, 30, GREEN_DIM, bold=True, new=False)
        tf = box(slide, cx + 0.18, y + 0.78, cw - 0.32, 1.5)
        write(tf, label, 11.5, TEXT, bold=True, spacing=1.18, font=MONO, new=False)
        if sub:
            write(tf, sub, 11, MUTED, before=6, spacing=1.24)

    block(slide, M, y + 2.78, CW, 1.24, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + 0.34, y + 3.0, CW - 0.68, 0.9)
    write(tf, "Two of these six are the ones that bite.", 15.5, TEXT, bold=True, new=False)
    write(
        tf,
        "Handlers registered after the listener starts miss the first events, and they are the events you most want "
        "to see. And media sent before SettingsApplied is thrown away, silently — which is why the bridge waits for "
        "the browser to say its speaker is live before it opens the Deepgram socket at all.",
        13,
        MUTED,
        before=5,
        spacing=1.3,
    )

    notes(
        slide,
        """
Reference slide. Leave it up while they work through Step 2.

Nobody writes this sequence -- it is in web/session.py and no step edits it. You
show it because every later step's bug has its root here, and because a room that
has seen the order stops guessing.

The two failure modes on the slide are the ones you will actually be asked about,
in Steps 3 and 4 rather than here:

Handlers after the listener: the greeting arrives within milliseconds of
SettingsApplied. Register late and you miss it, which reads as "the agent never
said anything."

Media before the handshake: silently discarded. In Step 4 that looks like an
agent that cannot hear you for the first second.

If anyone asks why the receive loop needs its own thread: websockets' send()
joins its receive thread with a ten-second close timeout as soon as a connection
starts closing. On the event loop that would freeze the whole server. That is
also why no coroutine ever calls send_media directly -- sends queue, and one
dedicated thread takes the risk.
        """,
    )


def slide_step3(prs: Presentation) -> None:
    """Step 3 gate: playback, and the queue between the network and the hardware.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "03",
        "Hear the agent",
        "0:45",
        "25 min",
        "LAB 1",
        "Route the agent's Flux TTS audio to your speaker, and hear the greeting out loud.",
        ["uv run steps/03-hear-the-agent/main.py"],
        "TODO 3.1 — one line, in the binary branch: player.send(audio)",
        "THE MENTAL MODEL",
        [
            "Audio arrives from the network in bursts — a few hundred milliseconds whenever Flux finishes a chunk.",
            "Your sound hardware consumes it at a constant rate, asking for exactly 128 samples at a time.",
            "A queue sits between the two. This step is that queue, and Step 5 is what happens when you need to "
            "throw it away.",
        ],
        "Everyone hears the greeting.",
    )
    notes(
        slide,
        """
First moment the workshop makes a sound. Worth waiting for the whole room.

Twenty-five minutes is generous for one line of code, and that is deliberate --
it absorbs the stragglers from Step 1 and Step 2.

Watch for: volume at zero, output routed to a disconnected Bluetooth device, and
people who did not press the button on the page. The page has to be clicked
before an AudioContext can be constructed, because browsers only unlock audio
inside a user gesture.

Check yourself, ask the room: why does the bridge wait for the browser to say its
speaker is live before opening the Deepgram socket? Answer: the greeting starts
arriving within milliseconds of SettingsApplied, and audio with nowhere to go is
audio thrown away.

When the room has heard it, name what just happened: three models, one socket,
and the reply was synthesised as it arrived rather than after a full sentence.
That is why it felt fast.
        """,
    )


def slide_step4(prs: Presentation) -> None:
    """Step 4 gate: microphone in, and the shortness of the outbound path.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "04",
        "Talk to the agent",
        "1:10",
        "30 min",
        "LAB 1",
        "Stream microphone audio to the agent and hold an actual conversation. This is the step where it becomes a "
        "voice agent.",
        ["uv run steps/04-talk-to-the-agent/main.py"],
        "TODO 4.1 — the entire outbound side: agent.send_media(audio)",
        "THE SHORTNESS IS THE LESSON",
        [
            "No voice activity detection. No silence trimming. No energy threshold. No “has the user finished” "
            "heuristic.",
            "Flux does turn detection inside the model, so there is nothing for client-side VAD to do. You forward "
            "bytes and get out of the way.",
        ],
        "Everyone holds a conversation. High point of the workshop — let it breathe.",
    )
    notes(
        slide,
        """
This is the high point. Do not rush it and do not talk over it.

When the room starts talking to their agents, stop facilitating for a minute and
let the noise happen. People will ask their agent something silly and laugh. That
is the moment they decide the workshop was worth attending.

Then draw the line under it: the outbound side of a voice agent is two lines,
because the model that already has the audio is the thing deciding when the turn
ended. Compare that to whatever VAD they have written before.

Check yourself, ask the room: why is there no voice activity detection anywhere
in this file? Answer: Flux does turn detection inside the model.

Watch for the agent interrupting itself. That is speakers feeding the microphone
past the echo canceller -- headphones, or turn the volume down. It is rarer than
it used to be, which makes it more confusing when it happens to one person.

Break is next, at 1:40. Before you send them off, say: anyone still behind should
open steps/05-barge-in/ during the break and start from there. Then use the ten
minutes to unstick people rather than to get coffee.
        """,
    )


def slide_break(prs: Presentation) -> None:
    """The ten-minute break, and what it is actually for.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)

    tf = box(slide, M, 1.5, 8.0, 1.9)
    write(tf, "1:40", 26, GREEN, bold=True, font=MONO, new=False)
    write(tf, "Break", 64, TEXT, bold=True, before=4, spacing=1.02)

    tf = box(slide, M, 3.7, 8.8, 1.2)
    write(tf, "Ten minutes.", 20, TEXT, bold=True, new=False)
    write(
        tf,
        "You have a working voice agent. After the break we break it — on purpose — and then fix it.",
        16,
        MUTED,
        before=6,
        spacing=1.3,
    )

    block(slide, M, 5.36, CW, 1.06, fill=AMBER_TINT, line=AMBER)
    pause_glyph(slide, M + 0.34, 5.75, 0.24)
    tf = box(slide, M + 1.02, 5.56, CW - 1.4, 0.68)
    write(tf, "Behind? Open steps/05-barge-in/ now.", 15, TEXT, bold=True, new=False)
    write(
        tf,
        "That folder already contains a complete, working Steps 1 through 4. Start there when we come back and you "
        "lose nothing but the typing.",
        13,
        AMBER,
        before=4,
        spacing=1.26,
    )

    notes(
        slide,
        """
Say the ten minutes and mean it. Coming back late compounds through the second
half, and Step 5 is the step you cannot afford to shorten.

Do not take the break yourself. This is your best window to unstick stragglers,
and the amber line on the slide is the thing to say to each of them
individually -- people will not take the shortcut on their own.

Worth a quick count while the room is empty: how many people finished Step 4?
That number tells you whether you are running the full path or cutting Step 8,
and it is the most useful thing to write down for next time.

If more than a couple of people are still on Step 2 or 3, plan to cut Step 8 and
Step 7's exercises now rather than discovering it at 3:05.
        """,
    )


def slide_step5(prs: Presentation) -> None:
    """Step 5 gate: feel the bug before fixing it.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "05",
        "Barge-in — run it first",
        "1:50",
        "25 min",
        "LAB 2",
        "Don't write anything yet. Interrupt it.",
        ["uv run steps/05-barge-in/main.py"],
        "Then: ask it something open-ended, wait for it to really get going, and talk over it.",
        "WHAT HAPPENS",
        [
            "It keeps going. Cheerfully, right across you, while the console prints >> UserStartedSpeaking.",
            "The server knew you had started talking. Your speaker did not care — it still had two seconds of "
            "speech in hand, and it played every byte of it.",
            "That gap between what the server knows and what your ears get is the whole step.",
        ],
        "Everyone experiences the bug before fixing it.",
        tag="do not cut",
    )
    notes(
        slide,
        """
The step people remember. Never cut it, and do not let anyone skip the first
part.

Run it together, as a room. Say: do not write anything yet. Ask your agent
something open-ended -- how a car engine works is reliable -- wait until it is
properly going, then talk over it.

The room will react. Someone will laugh. Wait for that before you say anything,
because the fix means nothing until they have felt the problem.

Then point at the console: >> UserStartedSpeaking is right there. The server knew.
This is not a detection failure, it is a playback failure, and it is entirely
yours to fix.

Check yourself, and it is the best question in the workshop: why isn't it enough
to tell the browser to flush? Answer: the pump would immediately refill the
browser's queue from the Python-side one. Clearing the far queue first and the
near queue second means the agent talks over the user a moment later rather than
immediately -- which is worse, because it looks like it nearly works.

The PortAudio equivalent, if someone is running --local: stop() drains the
buffer, playing everything already queued before stopping. abort() throws it away.
        """,
    )


def slide_two_queues(prs: Presentation) -> None:
    """Step 5: the two queues, and why order matters when you clear them.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Two queues sit between the agent and your ears",
        lede="Clearing only one leaves the bug in place — the pump simply refills it.",
    )

    stages = [
        ("Deepgram", "", None),
        ("queue in Python", "Outbox · web/audio.py", GREEN),
        ("queue in the browser", "PlaybackProcessor · worklets.js", GREEN),
        ("your ears", "", None),
    ]
    gap, n = 0.5, len(stages)
    cw = (CW - gap * (n - 1)) / n
    for i, (label, sub, color) in enumerate(stages):
        cx = M + i * (cw + gap)
        is_queue = color is not None
        block(slide, cx, y, cw, 1.46, fill=PANEL if is_queue else None, line=GREEN_DIM if is_queue else LINE)
        tf = box(slide, cx + 0.2, y + (0.34 if is_queue else 0.56), cw - 0.4, 0.9)
        write(tf, label, 15 if is_queue else 14, TEXT if is_queue else MUTED, bold=True, spacing=1.16, new=False)
        if sub:
            write(tf, sub, 10.5, MUTED, before=6, font=MONO, spacing=1.2)
        if i < n - 1:
            block(slide, cx + cw + 0.09, y + 0.71, gap - 0.18, 0.045, fill=GREEN_DIM, radius=None)

    tf = box(slide, M, y + 1.62, CW, 0.3)
    write(tf, "clear this one first  ↓", 10.5, AMBER, bold=True, caps=True, new=False)

    console(
        slide,
        M,
        y + 2.02,
        6.4,
        [
            ('elif message_type == "UserStartedSpeaking":', TEXT),
            ("    player.clear()", GREEN),
        ],
        size=13,
    )

    block(slide, 7.5, y + 2.02, W - 7.5 - M, 1.68, fill=PANEL_HI, line=LINE)
    tf = box(slide, 7.5 + 0.3, y + 2.2, W - 7.5 - M - 0.6, 1.34)
    write(tf, "clear() drops the Python side first,", 13.5, TEXT, bold=True, new=False)
    write(
        tf,
        "then tells the browser to flush. Both go through call_soon_threadsafe, which is FIFO, so that ordering "
        "holds. This is why Outbox is a deque and not a Queue — you cannot selectively drop from a queue.",
        12,
        MUTED,
        before=4,
        spacing=1.26,
    )

    footnote(
        slide,
        "Get the order backwards and the agent talks over the user a moment later instead of immediately. That is "
        "worse than not fixing it, because it looks like it nearly works.",
        y=y + 3.94,
        color=AMBER,
        size=13,
    )

    notes(
        slide,
        """
Show this after the room has heard the bug, never before.

The diagram is the whole explanation: two queues, and the one nearest Deepgram
has to go first. If you clear the browser and leave the Python side full, the
pump refills the browser within milliseconds.

Worth saying explicitly, because it is the transferable lesson: the failure mode
of clearing in the wrong order is not "it still does not work." It is "it works
in the demo you just did and fails when the answer is long." Those are the bugs
that ship.

The deque detail earns its place if anyone asks why not asyncio.Queue: you cannot
selectively drop items from a Queue, and Outbox.drop_audio has to remove audio
while leaving control frames alone.

After this, the workshop is finished as a piece of engineering. Everything from
Step 6 on is making it yours.
        """,
    )


# --- Steps 6 to 8: make it real ---------------------------------------------


def slide_step6(prs: Presentation) -> None:
    """Step 6 gate: prompt, voice, greeting, model.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "06",
        "Make it yours",
        "2:15",
        "20 min",
        "LAB 4",
        "The least code in the workshop and the most to play with. Four fields, and only one of them changes what "
        "the agent says.",
        ["uv run steps/06-make-it-yours/main.py"],
        "Change one thing at a time and run it again. That is the whole exercise.",
        "FOUR FIELDS",
        [
            "prompt — personality, job, boundaries. The standing instructions the model sees ahead of every turn.",
            "speak.provider.model — the voice. Purely how it sounds; it changes nothing about what the agent says.",
            "greeting — the opening line, and the only thing it says before you speak.",
            "think.provider.model — the brain, and its temperature.",
        ],
        "Go around the room and demo a few.",
        takehome="Take-home: Step 6b moves the brain to Amazon Bedrock, in your account. See the appendix.",
    )
    notes(
        slide,
        """
The most fun twenty minutes, and the easiest to under-run. Make the demo happen.

At the gate, go around and have four or five people put their agent on the room's
speakers. It costs three minutes and it is the only time people hear what
somebody else built.

The one instruction that matters: change one thing at a time. Swap the voice and
run the same conversation -- same words, different job applicant. Then raise the
temperature and watch it start improvising.

Voice catalogue is at developers.deepgram.com/docs/tts-models. A misspelled voice
warns rather than fails, so someone whose voice did not change should read the
console for >> Agent warning.

Check yourself, ask the room: name the two prompt instructions that matter for
speech but not for chat. Answer: tell it that it is speaking -- no markdown,
bullets, or emoji -- and tell it to be brief. Next slide covers both.

Then, in one sentence: there is an optional Step 6b that moves the LLM into your
own AWS account on Bedrock. It needs model access granted per model and per
region, which is calendar time rather than work, so it is a take-home unless this
room already has AWS. Do not spend more than a sentence on it.
        """,
    )


def slide_prompt_for_speech(prs: Presentation) -> None:
    """Step 6: the two prompt rules that are specific to voice.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Writing a prompt for speech",
        lede=(
            "Prompting a voice agent differs from prompting a chatbot in two specific ways, and both bite "
            "immediately."
        ),
    )

    rules = [
        (
            "Tell it that it's speaking.",
            "LLMs default to writing. Without an explicit instruction you get markdown, and text-to-speech reads it "
            "literally — **important** comes out as “star star important star star.”",
        ),
        (
            "Tell it to be brief.",
            "A four-sentence answer that scans fine in a chat window feels interminable when you have to sit "
            "through it. Voice punishes verbosity in a way text doesn't.",
        ),
    ]
    gap = 0.34
    cw = (CW - gap) / 2
    for i, (head, body) in enumerate(rules):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 1.86, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.34, y + 0.3, cw - 0.68, 1.4)
        write(tf, head, 18, GREEN, bold=True, spacing=1.14, new=False)
        write(tf, body, 13, MUTED, before=9, spacing=1.3)

    console(
        slide,
        M,
        y + 1.98,
        CW,
        [
            ("You are Sam, a barista at a small coffee shop. Your job is to take a", TEXT),
            ("drink order: the drink, the size, and the name for the cup.", TEXT),
            ("", TEXT),
            ("You are speaking out loud. Never use markdown, bullet points, lists,", GREEN),
            ("or emoji. Keep every reply to one or two sentences.", GREEN),
            ("", TEXT),
            ("Ask for one thing at a time, then read the order back and confirm it.", TEXT),
        ],
        size=12.5,
    )

    footnote(
        slide,
        "And keep it short for a third reason: every token goes to the model on every turn, so a long prompt slows "
        "down the first reply.",
        y=6.7,
        color=MUTED,
        size=13,
    )

    notes(
        slide,
        """
Two rules, thirty seconds, then let them write.

The markdown one always gets a laugh when you say it out loud in the star-star
voice. Use that -- it makes the rule stick, and it is the single most common
thing people get wrong on their first voice prompt.

The brevity rule is the one people resist, because a longer answer looks more
helpful in a text box. Point out that they will discover this themselves within
two exchanges of their own agent.

The barista prompt on the slide is a starting point, not the exercise. The
exercise is pushing at it: ask Sam for the weather, ask it to write you a Python
script, ask what it thinks of a film. A prompt that holds under pressure is the
difference between a demo and a product.

Worth mentioning for anyone thinking about production: prompt injection against a
voice agent is the same problem as against a chat agent, except your attacker is
talking out loud.
        """,
    )


def slide_step7(prs: Presentation) -> None:
    """Step 7 gate: function calling, in four hops.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "07",
        "Function calling",
        "2:35",
        "30 min",
        "LAB 5",
        "Ask it what time it is and it will confidently make something up, because an LLM has no clock. This step "
        "closes that gap and every gap like it.",
        ["uv run steps/07-function-calling/main.py"],
        "A phone banking agent for Contoso Bank that answers from your data, not its imagination.",
        "FOUR HOPS",
        [
            "1  You advertise a function — name, description, and a JSON Schema for its parameters, in SETTINGS.",
            "2  The LLM decides to call it, and Deepgram sends you a FunctionCallRequest with the arguments it chose.",
            "3  Your Python runs and returns a result.",
            "4  The model works that result into its reply, and the agent speaks it.",
        ],
        "Sketch a function for your own use case — the most useful two minutes in the workshop.",
        takehome="Take-home: Step 7b points the same machinery at a clinic. See the appendix.",
    )
    notes(
        slide,
        """
The step that turns a demo into something a business would pay for. Thirty
minutes, and the gate is the valuable part.

At the gate, do the sketching exercise properly: ask everyone to write down one
thing their agent would need to look up that a model cannot possibly know, then
have three people say theirs out loud. It reliably produces the best conversation
of the session, and it is where people work out whether they can actually use
this.

Before they start, plant the clock demo: ask your agent what time it is and
listen to it invent an answer with total confidence. Ten seconds, and it makes
the case for the whole step.

Check yourself, ask the room: what decides whether a function runs on your
machine or on Deepgram's? Answer: setting endpoint moves execution to Deepgram;
omitting it keeps the function client-side.

Next slide is the three things that will bite them. Show it before they start
rather than after they are stuck.

If you are running long, cut Step 7's exercises rather than the step. And in one
sentence at the end: Step 7b is the same machinery pointed at a clinic, it needs
nothing beyond the key you already have, and it teaches the two things banking
does not.
        """,
    )


def slide_step7_bites(prs: Presentation) -> None:
    """Step 7: the three mistakes worth pre-empting.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Three things that will bite you",
        lede="Show this before they start, not after they're stuck.",
    )

    numbered_rows(
        slide,
        y - 0.06,
        [
            (
                "1",
                "Omitting endpoint is what makes a function client-side",
                "Leave it unset and Deepgram sends the call down the socket to you. Set it and Deepgram calls your "
                "HTTP API itself, and you never see it. Both are valid — client-side is what you want when the "
                "function touches local state or credentials you'd rather not put on a public URL.",
            ),
            (
                "2",
                "The description is a prompt",
                "It is the only thing the LLM reads when deciding whether to call your function. “Get the current "
                "time” is weak. “Get the current date and time in a given IANA timezone. Use this whenever the user "
                "asks what time it is” is one the model can act on. A function that never fires usually has a "
                "description problem, not a wiring problem.",
            ),
            (
                "3",
                "on_message may not block",
                "It runs on the SDK's receive loop, which is also the thread carrying audio. A slow database call "
                "there does not just delay the reply, it stalls the conversation. That is why handle_function_call "
                "looks the way it does, and it holds on both the browser and --local paths.",
            ),
        ],
        row_h=1.32,
        gap=0.18,
        body_size=13,
        head_size=16,
    )

    footnote(
        slide,
        "The answer key for any of this is steps/99-final/main.py — the finished agent, every step applied.",
        y=6.74,
    )

    notes(
        slide,
        """
Reference slide. Leave it up for the whole step.

Number two is the one that eats time. When somebody says "it just isn't calling
my function," the description is the first thing to read, not the schema and not
the wiring. Say that out loud before they start and you will save several
one-on-one debugging sessions.

Number three is the one with production consequences. The threading rule the
whole design rests on: on_message runs on the receive loop, which is also
carrying audio, so it may not block. People will copy this pattern into real
systems, and a synchronous database call in that handler is a stall in the
conversation rather than a slow reply.

Number one is a design decision rather than a bug, but it is worth being precise
about because both directions are legitimate. Client-side for local state and
credentials; server-side when the function is just an HTTP call and you would
rather Deepgram made it.

If someone finishes early, point them at the transfer_funds exercise -- a
state-changing call needs confirmation before it fires, and a prompt is not an
authorization layer.
        """,
    )


def slide_step8(prs: Presentation) -> None:
    """Step 8 gate: the dials, and the trade they represent.

    Args:
        prs: The presentation to add to.
    """
    slide = gate_slide(
        prs,
        "08",
        "Optimization",
        "3:05",
        "15 min",
        "LAB 3",
        "Your agent is finished. Everything from here is how it feels — and this is what decides whether people "
        "enjoy talking to what you built.",
        ["uv run steps/08-optimize/main.py"],
        "Dials rather than code, which is what makes it the clean take-home.",
        "THREE DIALS",
        [
            "eot_threshold — 0.5 to 0.9. How much confidence Flux needs before calling the turn over. Raise it and "
            "it stops interrupting people who pause to think.",
            "eot_timeout_ms — ends the turn after this much silence regardless of confidence.",
            "eager_eot_threshold — at or below eot_threshold. Starts the LLM on a probable turn end and discards "
            "the work if the user keeps talking. Lower latency, more LLM calls.",
        ],
        "No gate. Stragglers catch up here, and this step sheds time cleanly.",
        tag="cuttable  ·  good take-home",
        tag_color=AMBER,
    )
    notes(
        slide,
        """
Pace-recovery step. It exists so you can end on time.

If you are running long, this is the first thing to cut, and it cuts cleanly --
it is configuration rather than code, the LAB.md reads fine unattended, and
nothing downstream depends on it. Say "this one is your homework" and go to the
wrap.

If you have the time, it is a good step, because it is the first time the room
has to make a judgement call rather than follow an instruction.

The sentence that carries it: no value is correct in the abstract. A drive-through
agent wants speed. An agent taking a credit card number over the phone wants
patience. You are choosing which mistake to make.

Check yourself, ask the room: your agent keeps cutting people off when they pause
to think. Which dial, and which way? Answer: raise eot_threshold -- it demands
more confidence before Flux calls the turn over.

Note there is no amber gate on this slide. Deliberate. Stragglers are catching up
during it, so holding the room here would defeat the point.
        """,
    )


def slide_find_your_setting(prs: Presentation) -> None:
    """Step 8: the four-run experiment.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Find your setting",
        lede="Run it four times, holding the same short conversation, and watch the latency line alongside.",
    )

    headers = ["EOT_THRESHOLD", "EOT_TIMEOUT_MS", "WHAT TO LISTEN FOR"]
    rows = [
        ("0.5", "5000", "Snappy, and it interrupts you when you pause", False),
        ("0.9", "5000", "Patient to the point of feeling slow", False),
        ("0.7", "500", "Cuts you off on any real pause", False),
        ("0.7", "5000", "Where you started — the balance you've been using all afternoon", True),
    ]
    widths = [2.2, 2.2, CW - 4.4 - 0.44]
    rh = 0.72

    x = M
    for header, cwidth in zip(headers, widths, strict=True):
        tf = box(slide, x + 0.24, y, cwidth - 0.24, 0.3)
        write(tf, header, 10.5, GREEN, bold=True, caps=True, font=MONO, new=False)
        x += cwidth + 0.22

    for i, (thresh, timeout, listen, highlight) in enumerate(rows):
        ry = y + 0.42 + i * (rh + 0.14)
        block(
            slide,
            M,
            ry,
            CW,
            rh,
            fill=GREEN_TINT if highlight else PANEL,
            line=GREEN_DIM if highlight else LINE,
        )
        x = M
        for value, cwidth, mono in zip((thresh, timeout, listen), widths, (True, True, False), strict=True):
            tf = box(slide, x + 0.24, ry + rh / 2 - 0.14, cwidth - 0.3, 0.32)
            write(
                tf,
                value,
                15 if mono else 13.5,
                (GREEN if highlight else TEXT) if mono else MUTED,
                bold=mono,
                font=MONO if mono else BODY,
                new=False,
            )
            x += cwidth + 0.22

    block(slide, M, 6.06, CW, 0.94, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + 0.34, 6.26, CW - 0.68, 0.6)
    write(tf, "No value is correct in the abstract.", 15, TEXT, bold=True, new=False)
    write(
        tf,
        "A drive-through agent wants speed; an agent taking a credit card number over the phone wants patience. "
        "You're choosing which mistake to make.",
        13,
        MUTED,
        before=4,
    )

    notes(
        slide,
        """
Leave it up while they run the four passes.

Insist on the same sentence every time, out loud, with a real hesitation in it:
"I'd like to order a... uh... large coffee." Four different settings, one
sentence, and the difference is unmistakable. People who improvise a new sentence
each run learn nothing.

Row three is the instructive one. A 500 ms timeout overrides a perfectly good
confidence score, which is how you discover that the timeout is a backstop and
not a tuning dial.

The latency line in the console is the honest half of this step. Patience costs
milliseconds, and the room should see the number rather than take your word for
it.

Going further, if anyone finishes early: eager_eot_threshold, 0.3 to 0.9, at or
below eot_threshold. It starts the LLM on a probable turn end and throws the work
away if the user keeps talking -- lower latency, more LLM calls. That trade is
the most production-relevant thing in this step.
        """,
    )


# --- Wrap --------------------------------------------------------------------


def slide_what_you_built(prs: Presentation) -> None:
    """Wrap: the inventory, and what deliberately isn't in it.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "What you built",
        kicker="wrap  ·  3:20",
        lede="Roughly 350 lines, most of them comments, doing six things.",
    )

    items = [
        "Opens a WebSocket and negotiates an agent configuration",
        "Streams microphone audio without blocking the thread that plays the reply",
        "Plays the agent's speech as it arrives, rather than after a full sentence",
        "Cuts playback the instant you interrupt",
        "Calls your Python and speaks the result",
        "Takes its turns on terms you chose, and measures what they cost",
    ]
    gap = 0.34
    cw = (CW - gap) / 2
    for i, item in enumerate(items):
        col, idx = (0, i) if i < 3 else (1, i - 3)
        ix = M + col * (cw + gap)
        iy = y + idx * 0.88
        block(slide, ix, iy, cw, 0.76, fill=PANEL, line=LINE)
        tick = block(slide, ix + 0.24, iy + 0.26, 0.24, 0.24, fill=None, line=GREEN, radius=0.5)
        tick.shadow.inherit = False
        tf = box(slide, ix + 0.24, iy + 0.29, 0.24, 0.2)
        write(tf, "✓", 10, GREEN, bold=True, align=PP_ALIGN.CENTER, new=False)
        tf = box(slide, ix + 0.62, iy + 0.19, cw - 0.86, 0.52)
        write(tf, item, 13, TEXT, spacing=1.22, new=False)

    block(slide, M, y + 2.86, CW, 1.1, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + 0.34, y + 3.06, CW - 0.68, 0.72)
    write(tf, "And notice what isn't in it.", 15.5, GREEN, bold=True, new=False)
    write(
        tf,
        "No device handling, no resampling, no chunking, no permission prompts, and no voice activity detection. "
        "The audio layer is all in web/, which every step shares and nobody edits — and the fact that it is "
        "replaceable is the point of keeping it out of the step files.",
        13,
        MUTED,
        before=5,
        spacing=1.28,
    )

    notes(
        slide,
        """
Sixty seconds. It is a victory lap, and the room has earned it.

Read the six lines, or better, ask the room to tell you what the agent does now.
People underestimate what they built because each step felt small.

The closing point is the last panel: none of the hard audio work is in their file.
No resampling, no chunking, no device handling, no VAD. That is a design choice
they can copy -- keep the transport out of the agent, and swapping a browser for
a Twilio media stream never touches agent code.

Point at steps/99-final/README.md for the extension ideas, and say the finished
agent is there as a reference implementation when one of their steps misbehaves
later.

Then leave real room for questions. The next two slides are for the people who
want somewhere to go, and the links slide is the one to leave up as the room
empties out.
        """,
    )


def slide_take_homes(prs: Presentation) -> None:
    """Wrap: the two optional steps, and why neither is in the run of show.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Two take-homes",
        lede="Detours, not links in the chain. Skip either and the next numbered step continues exactly as it would.",
    )

    take_homes = [
        (
            "STEP 6b  ·  15 MIN",
            "Bring your own LLM",
            "Moves the agent's brain to Amazon Bedrock, in your own AWS account. Deepgram makes the call as you, "
            "with credentials you hand it, against an endpoint you name. Your account, your model access, your bill.",
            "Needs Bedrock model access, granted per model and per region — which for a fresh account is calendar "
            "time rather than work. That's why it isn't in today's run.",
            AMBER,
        ),
        (
            "STEP 7b  ·  15 MIN",
            "A second vertical, in healthcare",
            "A clinic scheduling agent instead of a bank. Same function-calling machinery, pointed at a domain where "
            "the hard problem is what the agent must never repeat.",
            "Needs nothing beyond the Deepgram key you already have. Teaches the two things banking doesn't: "
            "keyterms, and holding back data the agent must not say aloud.",
            GREEN,
        ),
    ]
    gap = 0.34
    cw = (CW - gap) / 2
    for i, (kicker, title, body, caveat, color) in enumerate(take_homes):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 3.5, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.34, y + 0.3, cw - 0.68, 2.9)
        write(tf, kicker, 10.5, color, bold=True, caps=True, new=False)
        write(tf, title, 21, TEXT, bold=True, before=8, spacing=1.14)
        write(tf, body, 13, MUTED, before=10, spacing=1.3)
        write(tf, caveat, 12.5, color, before=12, spacing=1.28)

    footnote(
        slide,
        "Both read fine unattended — each LAB.md carries its own answer key, because nothing downstream can be one. "
        "Appendix has a slide on each.",
        y=6.5,
        color=MUTED,
        size=13,
    )

    notes(
        slide,
        """
Thirty seconds, and only if the room has energy left. Otherwise say one sentence
each and go to the links.

The honest framing on 6b: the step takes fifteen minutes and getting AWS Bedrock
model access takes days. If anyone wants it, tell them to request access before
they sit down to do it, not after. Access is per model and per region, and some
families need a one-time use case form that is not instant.

7b is the better recommendation for almost everyone in the room, because it costs
nothing but time and it teaches the habit that survives into production: return
only what the agent may say aloud. Filter at the boundary, not in the prompt.

If the room is a healthcare room, say you would swap 7b in for Step 7's exercises
rather than adding it on top.

Anyone who came from the Pipecat edition of this workshop and wants 6b: warn them
that AWS_BEARER_TOKEN_BEDROCK does nothing here. It is a botocore convenience and
there is no botocore in this stack.
        """,
    )


def slide_where_next(prs: Presentation) -> None:
    """Wrap: five directions out of the finished agent.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Where to go from here",
        lede="Five things you can try tonight, in rising order of how much they change.",
    )

    numbered_rows(
        slide,
        y - 0.1,
        [
            (
                "1",
                "Keyterms",
                "Bias recognition toward your product names, SKUs, and jargon Flux would otherwise mishear. One "
                "line of configuration, and usually the highest-leverage accuracy fix for a domain-specific agent.",
            ),
            (
                "2",
                "Multilingual",
                "Change the listen model to flux-general-multi for automatic language detection, and pass "
                "language_hints when you know the likely languages.",
            ),
            (
                "3",
                "Eager end-of-turn",
                "Add eager_eot_threshold to start the LLM on a probable turn end and discard the work if the user "
                "keeps talking. Lower latency, more LLM calls.",
            ),
            (
                "4",
                "Mid-conversation updates",
                "send_update_prompt changes instructions without reconnecting. send_inject_agent_message makes the "
                "agent say something unprompted, and send_inject_user_message makes function calling testable "
                "without a microphone.",
            ),
            (
                "5",
                "Telephony",
                "Drop SAMPLE_RATE to 8000 and switch encoding to mulaw to match what phone networks carry. main.py "
                "is otherwise unchanged — replacing web/ with a Twilio media stream is the real work.",
            ),
        ],
        row_h=0.94,
        gap=0.12,
        body_size=12.5,
        head_size=15,
    )

    notes(
        slide,
        """
Leave it up during questions. Do not read five bullets to a room at 3:20.

If you only point at one, point at keyterms. It is one line, it is already wired
into the listen provider in the finished agent, and it is the change that most
often turns a demo that mishears product names into something a business will
tolerate.

Telephony is last on the slide because it is the one people ask about first and
it is the biggest job -- though the interesting part is that main.py does not
change. The work is replacing web/, which is exactly why the audio layer lives
outside the step files.

Mid-conversation updates is the sleeper. send_inject_user_message makes function
calling testable without a microphone, which is how you would write tests for any
of this.

Everything here is in steps/99-final/README.md, so nobody has to write it down.
        """,
    )


def slide_links(prs: Presentation) -> None:
    """Wrap: the links, sized to hold QR codes.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Keep going",
        lede="Bring the questions you didn't get to. Both community channels are staffed by people who work on this.",
    )

    links = [
        ("Sign up / API key", "console.deepgram.com/signup"),
        ("Voice Agent docs", "developers.deepgram.com/docs/voice-agent"),
        ("Flux docs", "developers.deepgram.com/docs/flux"),
        ("Voice catalogue", "developers.deepgram.com/docs/tts-models"),
        ("Discord", "discord.gg/xWRaCDBtW4"),
        ("GitHub Discussions", "github.com/orgs/deepgram/discussions"),
    ]
    gap = 0.24
    cw = (CW - gap * 2) / 3
    ch = 1.94
    for i, (label, url) in enumerate(links):
        col, row = i % 3, i // 3
        cx = M + col * (cw + gap)
        cy = y + row * (ch + 0.26)
        block(slide, cx, cy, cw, ch, fill=PANEL, line=LINE)
        qr = 1.16
        block(slide, cx + 0.26, cy + (ch - qr) / 2, qr, qr, fill=None, line=LINE, radius=0.03)
        tf = box(slide, cx + 0.26, cy + ch / 2 - 0.1, qr, 0.24)
        write(tf, "QR", 10, DIM, bold=True, align=PP_ALIGN.CENTER, new=False)
        tf = box(slide, cx + qr + 0.44, cy + 0.42, cw - qr - 0.7, 1.2)
        write(tf, label, 14, TEXT, bold=True, spacing=1.14, new=False)
        write(tf, url, 10.5, GREEN, before=6, font=MONO, spacing=1.2)

    notes(
        slide,
        """
Leave this up as the room empties. It is the last thing anyone photographs.

The six QR slots ship empty on purpose -- generate them from the links table at
the end of FACILITATOR.md and drop them in before a conference floor, where
nobody will type a URL. For an internal team session, the URLs alone are fine.

Say two things while it is up:

Invite people into Discord while they still have the agent running. A question
asked with a live terminal gets a better answer than one asked three days later
from memory.

Share the repo link again, and point at steps/99-final/README.md. That is where
the extension ideas live, and it is the file people actually come back to.

Then, for your own sake, write down three things before you leave: which step
consumed the most time, how many people finished Step 7, and which environment
failures actually occurred. That is what you will change before the next run.
        """,
    )


# --- Appendix ----------------------------------------------------------------


def slide_appendix_divider(prs: Presentation) -> None:
    """Appendix divider.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)

    tf = box(slide, M, 2.1, 9.0, 1.8)
    write(tf, "appendix", 12, GREEN, bold=True, caps=True, new=False)
    write(tf, "Backup slides", 50, TEXT, bold=True, before=8, spacing=1.04)

    tf = box(slide, M, 4.06, 8.6, 0.7)
    write(
        tf,
        "Not part of the run of show. Jump to them when the room needs them.",
        16,
        MUTED,
        spacing=1.3,
        new=False,
    )

    contents = [
        "The five failures you'll actually hit",
        "Check-yourself answers",
        "Step 6b — bring your own LLM",
        "Step 7b — a second vertical",
        "Running in another region",
    ]
    gap, n = 0.16, len(contents)
    cw = (CW - gap * (n - 1)) / n
    for i, item in enumerate(contents):
        cx = M + i * (cw + gap)
        block(slide, cx, 5.2, cw, 1.0, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.2, 5.4, cw - 0.4, 0.7)
        write(tf, item, 12, MUTED, spacing=1.24, new=False)

    notes(
        slide,
        """
Never presented. This slide exists so you can find the backup material quickly
in presenter view without scrolling past thirty slides.

Practical tip: note these slide numbers on a sticky before you start. The failures
slide is the one you will want inside the first twenty minutes, and hunting for it
in front of a room while somebody's microphone does not work is worse than
describing it from memory.

If you are running the deck in presenter mode, PowerPoint's slide-jump shortcut is
typing the slide number and pressing Enter.
        """,
    )


def slide_failures(prs: Presentation) -> None:
    """Appendix: the environment failures that actually happen.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "The five failures you'll actually hit",
        kicker="appendix",
        lede="In order of frequency. Four of the five surface during Step 1, which is what Step 1 is for.",
    )

    numbered_rows(
        slide,
        y - 0.24,
        [
            (
                "1",
                "Microphone permission",
                "The prompt appears during Step 1 and people click past it. Icon at the left of the address bar, "
                "allow, press the button again. On macOS also check System Settings → Privacy & Security → "
                "Microphone and confirm the browser is enabled — it's the browser now, not the terminal.",
            ),
            (
                "2",
                "The page is open on the wrong address",
                "getUserMedia and AudioWorklet need a secure context, which browsers grant to localhost and nothing "
                "else, so a LAN address makes the audio API silently vanish. Symptom: the Connect button stays "
                "disabled with a red box above it. Use the http://127.0.0.1:8000 the terminal printed.",
            ),
            (
                "3",
                "Bluetooth headsets",
                "Activating the microphone flips many headsets into a low-quality mono profile, and some then fail "
                "to open at all. Reconnect, or switch to wired. This is why the pre-event email asks for wired.",
            ),
            (
                "4",
                "The agent talks to itself",
                "Laptop speakers feeding the laptop microphone, past the echo canceller. Much rarer than it used to "
                "be, which makes it more confusing when it happens to one person. Headphones, or turn the volume "
                "down. Step 1 reports whether the browser actually granted echo cancellation.",
            ),
            (
                "5",
                "A truncated API key",
                "Copy-paste drops characters or adds a space. Step 1 makes a real authenticated call rather than "
                "checking it is non-empty, so this surfaces at once as “Deepgram rejected the key.”",
            ),
        ],
        row_h=0.74,
        gap=0.02,
        body_size=11,
        head_size=13,
    )

    block(slide, M, 5.92, CW, 0.96, fill=AMBER_TINT, line=AMBER)
    tf = box(slide, M + 0.3, 6.09, CW - 0.6, 0.7)
    write(tf, "And one that isn't an environment failure at all:", 12.5, AMBER, bold=True, new=False)
    write(
        tf,
        "“The agent refused these settings.” The machine is fine — a model isn't available where that person is "
        "connecting, and Step 1 prints which. Same fix for the whole room: DEEPGRAM_REGION=global, or change the "
        "named model everywhere.",
        11.5,
        MUTED,
        before=3,
        spacing=1.24,
    )

    notes(
        slide,
        """
Your most-used slide. Know where it is before you start.

Numbers one and two account for most of what you will see, and both are one-line
fixes once you know to look. Number two is the sneaky one, because the audio API
does not throw an error a person could search for -- ctx.audioWorklet is simply
undefined. If someone is on a LAN address, nothing they read online will help them.

Number six, only if you are running Step 6b: Bedrock model access. Not on the main
line, but by far the longest to fix, because you cannot fix it in the room. Two
cheaper mistakes shadow it and are worth ruling out first -- a region mismatch
between the credentials and the endpoint URL, and an IAM user holding
bedrock:InvokeModel but not bedrock:InvokeModelWithResponseStream. The agent
streams, so the non-streaming permission alone is not enough.

For the amber panel: the room shares a .env layout, so if one person hits a
refused model, everyone connecting to that region will. Fix it once, out loud, for
the whole room.
        """,
    )


def slide_answers_core(prs: Presentation) -> None:
    """Appendix: check-yourself answers for Steps 0 to 5.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Check-yourself answers  ·  Steps 0 – 5",
        kicker="appendix",
        lede="The LAB.md files pose these. Here are the answers, for when you ask the room.",
    )

    answers = [
        (
            "0",
            "STT, LLM, TTS — and the LLM holds the conversation history.  ·  Turn-taking, provider handoffs, "
            "buffering, interruption signalling.  ·  Turn detection happens inside the Flux model, server-side; the "
            "client is still responsible for clearing queued playback on barge-in.  ·  The next folder.",
        ),
        ("1", "One. Deepgram brokers the OpenAI call, so your Deepgram key covers the LLM."),
        (
            "2",
            "listen configures speech-to-text; the LLM is named under think.  ·  The agent discards any media that "
            "arrives before the handshake completes, so audio sent early is silently lost.",
        ),
        (
            "3",
            "The greeting starts arriving within milliseconds of SettingsApplied, and audio with nowhere to go is "
            "audio thrown away. The bridge waits for the browser's start message before opening the Deepgram socket "
            "for exactly this reason.",
        ),
        ("4", "Flux does turn detection inside the model, so there's nothing for client-side VAD to do."),
        (
            "5",
            "Because the pump would immediately refill the browser's queue from the Python-side one. Clearing the "
            "far queue first and the near queue second means the agent talks over the user a moment later rather "
            "than immediately — which is worse, because it looks like it nearly works.",
        ),
    ]

    ry = y - 0.18
    for num, text in answers:
        lines = max(1, -(-len(text) // 145))
        h = 0.3 + lines * 0.215
        block(slide, M, ry, CW, h, fill=PANEL, line=LINE)
        tf = box(slide, M + 0.26, ry + 0.16, 0.4, 0.3)
        write(tf, num, 15, GREEN, bold=True, font=MONO, new=False)
        tf = box(slide, M + 0.82, ry + 0.15, CW - 1.1, h - 0.3)
        write(tf, text, 12, MUTED, spacing=1.26, new=False)
        ry += h + 0.12

    notes(
        slide,
        """
For you, not for the room. Do not project it before asking the question.

The two worth asking out loud, because the answers change how people build:

Step 4's -- why is there no VAD anywhere in this file. It is the moment the Flux
argument lands for people who have written voice activity detection before.

Step 5's -- why isn't it enough to tell the browser to flush. Nobody gets this
right first time, and the wrong answer is instructive: clearing the near queue
first produces a bug that looks almost fixed, which is the worst kind.

Step 0's four answers are quick fire. Use them to fill the last minute of the
overview block if you are running ahead.
        """,
    )


def slide_answers_rest(prs: Presentation) -> None:
    """Appendix: check-yourself answers for Steps 6 to 8 and the detours.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Check-yourself answers  ·  Steps 6 – 8",
        kicker="appendix",
        lede="Including the two optional steps, which carry their own answer keys because nothing downstream can.",
    )

    answers = [
        ("6", "Tell it that it's speaking — no markdown, bullets, or emoji — and tell it to be brief.", GREEN),
        (
            "6b",
            "Deepgram brokers OpenAI, so a model name is the whole configuration. It doesn't broker Bedrock: it "
            "makes the call as you, which needs credentials to make it with (provider.credentials) and an address "
            "to make it to (think.endpoint).",
            AMBER,
        ),
        ("7", "Setting endpoint moves execution to Deepgram; omitting it keeps the function client-side.", GREEN),
        (
            "7b",
            "Because the prompt is a request and the payload is a guarantee. Prose asking the model not to repeat a "
            "phone number can be talked around; a payload that never carried the phone number cannot.",
            AMBER,
        ),
        ("8", "Raise eot_threshold. It demands more confidence before Flux calls the turn over.", GREEN),
    ]

    ry = y - 0.14
    for num, text, color in answers:
        lines = max(1, -(-len(text) // 140))
        h = 0.32 + lines * 0.225
        block(slide, M, ry, CW, h, fill=PANEL, line=LINE)
        tf = box(slide, M + 0.26, ry + 0.19, 0.56, 0.3)
        write(tf, num, 15, color, bold=True, font=MONO, new=False)
        tf = box(slide, M + 0.98, ry + 0.18, CW - 1.26, h - 0.34)
        write(tf, text, 12.5, MUTED, spacing=1.28, new=False)
        ry += h + 0.12

    block(slide, M, ry + 0.12, CW, 1.04, fill=PANEL_HI, line=LINE)
    tf = box(slide, M + 0.34, ry + 0.28, CW - 0.68, 0.8)
    write(tf, "7b's answer is the one worth labouring.", 14, GREEN, bold=True, new=False)
    write(
        tf,
        "It generalises past healthcare and past voice: filter at the boundary, because the model may say anything "
        "you hand it — under pressure, or just because it was being helpful.",
        12.5,
        MUTED,
        before=4,
        spacing=1.26,
    )

    notes(
        slide,
        """
Same as the previous slide: yours, not the room's.

Step 6's answer is the one that saves people the most time on their own projects,
and it is worth asking even if you are behind. Two instructions, ten seconds.

Step 6b's answer is the useful mental model for the whole provider list. Brokered
means Deepgram holds the account and a model name is the entire configuration.
Bring-your-own means Deepgram calls as you, which needs credentials and an
address. Once someone has that distinction, Groq and Bedrock stop being special
cases.

Step 7b's is the one to labour if the room includes anyone building for a
regulated industry. A prompt is a request. A payload is a guarantee. Everything
else in a compliance conversation is downstream of understanding that difference.
        """,
    )


def slide_step6b(prs: Presentation) -> None:
    """Appendix: Step 6b, brokered versus bring-your-own providers.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Step 6b — Bring your own LLM",
        kicker="appendix  ·  optional  ·  15 min",
        lede="Move the agent's brain onto Amazon Bedrock, in your own AWS account. The provider list splits in two.",
    )

    columns = [
        (
            "BROKERED BY DEEPGRAM",
            GREEN,
            "OpenAI · Anthropic · Google · NVIDIA",
            "Deepgram holds the account, makes the call, and bills you. You name a model and you're done — which is "
            "why Step 6 let you swap gpt-4o-mini for gpt-4o without so much as an OpenAI account.",
        ),
        (
            "BRING YOUR OWN",
            AMBER,
            "Groq · AWS Bedrock",
            "Deepgram makes the call as you, with credentials you hand it, against an endpoint you name. Your "
            "account, your model access, your bill. This is the answer when the model has to run somewhere you "
            "control.",
        ),
    ]
    gap = 0.34
    cw = (CW - gap) / 2
    for i, (label, color, models, body) in enumerate(columns):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 1.84, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.32, y + 0.22, cw - 0.64, 1.42)
        write(tf, label, 11, color, bold=True, caps=True, new=False)
        write(tf, models, 14.5, TEXT, bold=True, before=6, font=MONO, spacing=1.16)
        write(tf, body, 12.5, MUTED, before=8, spacing=1.3)

    cards(
        slide,
        y + 2.06,
        1.34,
        [
            (
                "think.provider",
                "Which model, and the credentials to reach it. zai.glm-4.7-flash is the workshop default — fast and "
                "cheap, which matters more than raw capability when someone is waiting to hear a reply. Takes "
                "long-lived iam keys or short-lived sts ones, which additionally carry a session_token.",
            ),
            (
                "think.endpoint",
                "https://bedrock-runtime.{region}.amazonaws.com/ — and the region has to match your credentials'. "
                "Bedrock needs both settings; miss either and the handshake fails.",
            ),
        ],
        label_size=14,
        body_size=12,
        label_font=MONO,
    )

    block(slide, M, y + 3.56, CW, 1.04, fill=AMBER_TINT, line=AMBER)
    tf = box(slide, M + 0.34, y + 3.73, CW - 0.68, 0.78)
    write(tf, "Be clear-eyed about where the credential travels.", 14, AMBER, bold=True, new=False)
    write(
        tf,
        "Your AWS access key and secret go into the Settings message, over the WebSocket, to Deepgram. Scope an IAM "
        "user to bedrock:InvokeModelWithResponseStream on the one model ARN in use — the blast radius of a workshop "
        "credential should be one model in one region.",
        12,
        MUTED,
        before=4,
        spacing=1.26,
    )

    notes(
        slide,
        """
Backup slide. Use it if someone asks about lock-in, data residency, or running
their own model -- which someone usually does around Step 6.

The setting worth remembering is think.endpoint, and it is not Bedrock-specific.
Point it at anything that speaks the OpenAI Chat Completions format -- a
self-hosted model, a gateway in front of your own inference, a router, a proxy
that logs every completion -- and the agent talks to it. Bedrock is just the case
with enough structure that Deepgram gave it a provider type of its own.

Two practical warnings if anyone is going to actually do this:

Model access is per model and per region, and some families need a one-time use
case form that is not instant. Request it days ahead.

bedrock:InvokeModel is not enough. The agent streams, so it needs
bedrock:InvokeModelWithResponseStream.

Worth naming the graceful-degradation design too: think_settings() guards on the
presence of AWS credentials and falls back to the brokered provider without them.
The person whose model access never got approved can still run your file.

And the bill genuinely moves -- after 6b the request shows up in AWS, Bedrock,
Usage, and nowhere in the Deepgram console. Speech-to-text and text-to-speech are
still Deepgram's; only the middle of the pipeline changed hands.
        """,
    )


def slide_step7b(prs: Presentation) -> None:
    """Appendix: Step 7b, keyterms and the payload as enforcement.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Step 7b — A second vertical, in healthcare",
        kicker="appendix  ·  optional  ·  15 min",
        lede="Same machinery, a clinic instead of a bank. Two lessons banking never asks for.",
    )

    columns = [
        (
            "KEYTERMS, HEARD NOT EXPLAINED",
            "Run it once with the list empty and say: “I'm calling about my semaglutide prescription with Doctor "
            "Bergstrom.” Read what lands in the transcript — it won't be that sentence. Then uncomment the list and "
            "say it again. That's the entire feature, and it's one line.",
        ),
        (
            "THE PAYLOAD IS THE ENFORCEMENT LAYER",
            "return patient works, and hands a full date of birth, patient ID, and phone number to a language "
            "model — trusting a paragraph of prose to stop it repeating them. Return three keys instead. The prompt "
            "says the agent must not speak an ID; the payload makes it impossible.",
        ),
    ]
    gap = 0.34
    cw = (CW - gap) / 2
    for i, (label, body) in enumerate(columns):
        cx = M + i * (cw + gap)
        block(slide, cx, y, cw, 2.2, fill=PANEL, line=LINE)
        tf = box(slide, cx + 0.32, y + 0.28, cw - 0.64, 1.7)
        write(tf, label, 11.5, GREEN, bold=True, caps=True, new=False)
        write(tf, body, 12.5, MUTED, before=9, spacing=1.3)

    console(
        slide,
        M,
        y + 2.52,
        6.4,
        [
            ("return {", TEXT),
            ('    "verified": True,', GREEN),
            ('    "first_name": patient["first_name"],', GREEN),
            ('    "next_appointment": patient["next_appointment"],', GREEN),
            ("}", TEXT),
        ],
        size=12.5,
    )

    block(slide, 7.5, y + 2.52, W - 7.5 - M, 1.66, fill=GREEN_TINT, line=GREEN_DIM)
    tf = box(slide, 7.5 + 0.3, y + 2.72, W - 7.5 - M - 0.6, 1.3)
    write(tf, "A compliance control, written as a return statement.", 13.5, GREEN, bold=True, spacing=1.16, new=False)
    write(
        tf,
        "Minimum necessary — §164.502(b), the rule that you disclose only what the task requires — is precisely what "
        "this implements. And it costs nothing: you didn't lose a feature, you removed data the conversation never "
        "needed.",
        11.5,
        MUTED,
        before=5,
        spacing=1.26,
    )

    footnote(
        slide,
        "Demo only. Not a HIPAA-compliant system, and every record in it is invented. The LAB.md documents the "
        "distance between this file and something allowed near a real patient — and most of it isn't code.",
        y=6.6,
        color=CORAL,
        size=12,
    )

    notes(
        slide,
        """
Backup slide, and the one to reach for when someone asks about a regulated
industry or about stopping an agent leaking data.

Keyterms is best demonstrated rather than described. If you have two minutes and
a microphone, run it with the list empty and say the semaglutide sentence out
loud. The transcript comes back with a drug that does not exist, and the room
understands keyterms permanently. It is the highest-leverage accuracy fix
available for a domain-specific agent and it is one line of configuration.

The payload lesson is the transferable one, and it generalises well past
healthcare. The prompt is a request; the payload is a guarantee. Prose asking a
model not to repeat a phone number can be talked around -- under pressure, under
a clever question, or just because the model was being helpful. A payload that
never carried the phone number cannot.

If someone wants the compliance conversation: BAAs with every vendor that touches
PHI, including logging and observability, which is the one teams forget. Audit
controls that are not print statements. Retention and redaction decisions for the
audio as well as the transcript. And a documented risk analysis, training, and
breach process around all of it. There is no such thing as a HIPAA-certified
codebase -- it is a program you maintain, and it wants a lawyer involved early.

Deepgram is a Business Associate under HIPAA and offers BAAs to qualifying
Covered Entities on request. That is an account-team conversation, not a
configuration flag.
        """,
    )


def slide_regions(prs: Presentation) -> None:
    """Appendix: one line moves every step to another region.

    Args:
        prs: The presentation to add to.
    """
    slide = new_slide(prs)
    y = heading(
        slide,
        "Running in another region",
        kicker="appendix",
        lede="Global is the default. The EU and AU endpoints process audio inside those geographies.",
    )

    console(slide, M, y, 6.4, [("DEEPGRAM_REGION=eu    # global (default), eu, or au", GREEN)], size=13.5)

    tf = box(slide, M, y + 0.86, 6.4, 1.5)
    write(
        tf,
        "That is the whole change. Your key works in every region, no step's code names one, and every step reads "
        "the same .env — so one line moves Steps 1 through 8 together. The same mechanism points the workshop at a "
        "Deepgram Dedicated or self-hosted deployment.",
        13,
        MUTED,
        spacing=1.3,
        new=False,
    )

    block(slide, 7.5, y, W - 7.5 - M, 2.34, fill=PANEL, line=LINE)
    tf = box(slide, 7.5 + 0.32, y + 0.26, W - 7.5 - M - 0.64, 1.9)
    write(tf, "you don't have to test it yourself", 10.5, GREEN, bold=True, caps=True, new=False)
    write(
        tf,
        "Step 1 opens the same WebSocket Step 2 will, with the same three models, and reports whether the server "
        "accepted them.",
        12.5,
        TEXT,
        before=8,
        spacing=1.28,
    )
    write(
        tf,
        "Model availability is per-region and moves over time. A key that works and a host that answers still don't "
        "tell you a model is served there — and that is the question worth asking before a room of people hits "
        "Step 2.",
        12,
        MUTED,
        before=7,
        spacing=1.28,
    )

    block(slide, M, y + 2.62, CW, 1.0, fill=AMBER_TINT, line=AMBER)
    tf = box(slide, M + 0.34, y + 2.8, CW - 0.68, 0.68)
    write(tf, "As of August 2026, au needs one change.", 14, AMBER, bold=True, new=False)
    write(
        tf,
        "Flux TTS isn't served there yet, so flux-alexis-en is refused. Flux STT and the LLM are fine. Set speak to "
        "aura-2-thalia-en in every step's SETTINGS and the workshop runs unchanged — Step 6 teaches changing the "
        "voice anyway, so it costs the room nothing. eu and global both run it as shipped.",
        12,
        MUTED,
        before=4,
        spacing=1.26,
    )

    footnote(
        slide,
        "One asymmetry worth knowing: the Management API call Step 1 makes to prove your key works is global only. "
        "region.management_client() exists to keep that deliberate.",
        y=6.72,
    )

    notes(
        slide,
        """
Backup slide, and pre-flight reading if you are running this outside the global
endpoint.

Add the DEEPGRAM_REGION line to the pre-event email if attendees clone the repo
themselves. .env.example documents it, but people paste their key and stop
reading.

The important operational point is that you do not have to verify a region
yourself. Step 1 opens the real socket with the real three models and reports what
the server said, per attendee. That is stronger than any check you could run in
advance, and it means each person finds out on their own machine.

The au caveat is dated deliberately, because it will go stale. Step 1 tells each
attendee directly, so the room will find out even if this slide is out of date --
which is the right way round.

If someone asks about audio never leaving their network: that is self-hosted, and
on AWS it takes the form of Flux, Nova-3, and Aura-2 on SageMaker endpoints in
their own VPC. Speech models only -- SageMaker's network isolation blocks the
outbound LLM calls the Voice Agent orchestrator makes, so the Agent API itself
does not run there. Different architecture, not a setting.
        """,
    )


# --- Build -------------------------------------------------------------------

SLIDES = [
    slide_title,
    slide_what_youll_build,
    slide_prerequisites,
    slide_two_tracks,
    slide_run_of_show,
    slide_three_models,
    slide_architecture,
    slide_orchestration,
    slide_what_flux_changes,
    slide_one_job,
    slide_layout,
    slide_step1,
    slide_step1_working,
    slide_step2,
    slide_step2_ordering,
    slide_step3,
    slide_step4,
    slide_break,
    slide_step5,
    slide_two_queues,
    slide_step6,
    slide_prompt_for_speech,
    slide_step7,
    slide_step7_bites,
    slide_step8,
    slide_find_your_setting,
    slide_what_you_built,
    slide_take_homes,
    slide_where_next,
    slide_links,
    slide_appendix_divider,
    slide_failures,
    slide_answers_core,
    slide_answers_rest,
    slide_step6b,
    slide_step7b,
    slide_regions,
]


def main() -> None:
    """Build the deck and write it to slides/, overwriting what's there."""
    prs = Presentation()
    prs.slide_width = Emu(int(W * 914400))
    prs.slide_height = Emu(int(H * 914400))

    for builder in SLIDES:
        builder(prs)

    out = Path(__file__).resolve().parent.parent / "slides" / "deepgram-voice-agent-workshop.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote {out.relative_to(out.parent.parent)} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
